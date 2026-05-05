# build_briefs_v2.py
# ------------------------------------------------------------
# 후보기술개요서 자동 요약 봇 (통합 개선판)
# - 1-Stage LLM 호출 (Chat Completions API)
# - 논문/특허 자동 판별
# - 사용자 지침 기반 프롬프트 엔지니어링
# - macOS 호환
# ------------------------------------------------------------

import os
import re
import io

try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    pass
import json
import time
import argparse
from pathlib import Path
from copy import deepcopy
from dataclasses import dataclass, field
from typing import List, Dict, Any, Tuple, Optional

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from pptx.enum.text import PP_ALIGN

from openai import OpenAI

try:
    import google.generativeai as genai
    HAS_GEMINI = True
except ImportError:
    HAS_GEMINI = False

try:
    from PIL import Image
except Exception:
    Image = None

try:
    import requests
    from bs4 import BeautifulSoup
    HAS_REQUESTS = True
except Exception:
    HAS_REQUESTS = False

# API keys (env에서 자동 로드)
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "").strip()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "").strip()

# ============================================================
# Constants
# ============================================================
DEFAULT_MODEL = "gpt-5"
IS_GEMINI_MODEL = lambda m: m.lower().startswith("gemini")

# 새 템플릿 Shape 이름 (선택창 기반)
SHAPE_NUMBER_INFO = "number_info"
SHAPE_TITLE = "title"
SHAPE_SUMMARY = "summary_point_1-2"
SHAPE_DIAGRAM1 = "diagram1"
SHAPE_DIAGRAM2 = "diagram2"
SHAPE_CAPTION1 = "concept_diagram1"
SHAPE_CAPTION2 = "concept_diagram2"
SHAPE_SOURCE = "source_info"

# v2 템플릿: 15×1 표 기반 (기술목적/문제점/제안기술/개선효과)
TABLE_NAME = "표 10"
ROW_PURPOSE_HEADER = 0
ROW_PURPOSE_CONTENT = 1
ROW_PROBLEMS_HEADER = 2
ROW_PROBLEMS_CONTENT = 3
ROW_METHOD_HEADER = 4
# method title rows: 5, 7, 9, 11 / method detail rows: 6, 8, 10, 12
ROW_METHOD_TITLE_START = 5
ROW_METHOD_DETAIL_START = 6
ROW_EFFECT_HEADER = 13
ROW_EFFECT_CONTENT = 14

SLIDE_W: Optional[int] = None
SLIDE_H: Optional[int] = None

# ============================================================
# Timer utility (요청서 §9-3)
# ============================================================
def _timer_log(label: str, start: float) -> float:
    """단계별 소요시간 로그 출력. 현재 시각을 반환."""
    elapsed = time.time() - start
    print(f"[TIMER] {label}: {elapsed:.2f}s")
    return time.time()


# ============================================================
# Figure category classification (요청서 §3~§5)
# ============================================================
_METHOD_KEYWORDS = [
    "system", "structure", "circuit", "architecture", "proposed", "method",
    "framework", "model", "block diagram", "flowchart", "schematic",
    "control structure", "configuration", "concept", "mechanism", "design",
    "topology", "layout", "overview", "pipeline", "module", "flow", "diagram",
    # 한국어
    "시스템", "구조", "회로", "아키텍처", "제안", "프레임워크", "모델",
    "블록도", "흐름도", "개략도", "제어", "구성", "컨셉", "메커니즘", "설계",
]
_RESULT_KEYWORDS = [
    "result", "results", "experiment", "experimental", "test", "validation",
    "performance", "comparison", "efficiency", "error", "loss", "accuracy",
    "waveform", "simulation", "hil", "measurement", "evaluation", "improvement",
    "response", "spectrum", "bode", "histogram", "scatter", "bar chart",
    # 한국어
    "결과", "실험", "검증", "성능", "비교", "효율", "오차", "손실",
    "정확도", "파형", "시뮬레이션", "측정", "평가", "개선효과",
]


def _classify_figure_category(caption: str) -> Tuple[str, str, float]:
    """caption 키워드 기반 figure 카테고리 분류.

    Returns: (category, role_candidate, selection_score)
      - category: 'method_or_concept' | 'result_or_effect' | 'other'
      - role_candidate: 'top_image' | 'bottom_image' | ''
      - selection_score: 0.0~1.0 (키워드 매칭 밀도)
    """
    if not caption:
        return "other", "", 0.0
    cap_lower = caption.lower()
    method_score = sum(1 for kw in _METHOD_KEYWORDS if kw in cap_lower)
    result_score = sum(1 for kw in _RESULT_KEYWORDS if kw in cap_lower)
    total_kw = len(_METHOD_KEYWORDS) + len(_RESULT_KEYWORDS)
    max_score = max(method_score, result_score)
    density = min(1.0, max_score / max(total_kw * 0.05, 1.0))  # normalize

    if method_score > result_score:
        return "method_or_concept", "top_image", density
    elif result_score > method_score:
        return "result_or_effect", "bottom_image", density
    elif method_score > 0:  # tied but non-zero
        return "method_or_concept", "top_image", density
    return "other", "", 0.0


# ============================================================
# JSON Schema for unified 1-stage output
# ============================================================
BRIEF_SCHEMA: Dict[str, Any] = {
    "type": "object",
    "properties": {
        "doc_type": {"type": "string", "enum": ["paper", "patent"]},
        "title": {"type": "string"},
        "head_messages": {
            "type": "array", "items": {"type": "string"}, "minItems": 2, "maxItems": 2
        },
        "purpose": {
            "type": "array", "items": {"type": "string"}, "minItems": 2, "maxItems": 4
        },
        "prior_problems": {
            "type": "array", "items": {"type": "string"}, "minItems": 2, "maxItems": 4
        },
        "proposed_method": {
            "type": "array", "minItems": 3, "maxItems": 4,
            "items": {
                "type": "object",
                "properties": {
                    "title": {"type": "string"},
                    "details": {"type": "array", "items": {"type": "string"}, "minItems": 3, "maxItems": 3}
                },
                "required": ["title", "details"]
            }
        },
        "improvements": {
            "type": "array", "items": {"type": "string"}, "minItems": 2, "maxItems": 4
        },
        "conclusion": {
            "type": "array", "items": {"type": "string"}, "minItems": 2, "maxItems": 3
        },
        "paper_info": {
            "type": "object",
            "properties": {
                "journal_or_patent_office": {"type": "string"},
                "paper_title": {"type": "string"},
                "institution": {"type": "string"},
                "doi_or_patent_no": {"type": "string"},
                "year": {"type": "string"},
                "month": {"type": "string"}
            },
            "required": ["journal_or_patent_office", "paper_title", "institution", "doi_or_patent_no", "year", "month"]
        },
        "representative_figures": {
            "type": "array", "minItems": 2, "maxItems": 2,
            "items": {
                "type": "object",
                "properties": {
                    "page": {"type": "integer"},
                    "fig_number": {"type": "string"},
                    "fig_title": {"type": "string"},
                    "role": {"type": "string", "enum": ["concept", "result"]}
                },
                "required": ["page", "fig_number", "fig_title", "role"]
            }
        },
        "figure_captions_ko": {
            "type": "array", "items": {"type": "string"}, "minItems": 2, "maxItems": 2
        }
    },
    "required": [
        "doc_type", "title", "head_messages", "purpose", "prior_problems",
        "proposed_method", "improvements", "conclusion", "paper_info",
        "representative_figures", "figure_captions_ko"
    ]
}

# ============================================================
# System Prompt (핵심 — 웹 ChatGPT 품질 매칭)
# ============================================================
SYSTEM_PROMPT = """\
역할, 목표, 출력, 원칙

# 역할
너는 자동차 분야 10년차 연구원이고, 자동차 분야의 논문/특허의 기술적 핵심을 중심으로 요약 전문가이다. 
PDF 문서를 읽고, 아래 지침에 따라 한국어로 요약하여 지정된 JSON 형식으로 출력해

# 목표
논문 및 특허의 기술적 핵심을 중심으로 핵심 내용을 정해진 요약 관점(제목(기술명), 헤드메세지, 기술목적, 문제점, 제안기술, 개선효과)에 맞춰 정리해 

# 작성 원칙
## 공통 작성 규칙
- 약어·전문 용어는 첫 등장 시 괄호 안에 영어 원문 또는 풀네임을 병기 (예: "비가시 경로(NLOS, Non-Line-of-Sight)", "PPP-RTK(Precise Point Positioning-Real Time Kinematic)")
- 구체적인 알고리즘명·구조·프로세스를 포함하여 상세하게 작성. 단, 수식(수학식·방정식)은 그대로 나열하지 말고 해당 수식의 의미를 개념적으로 설명할 것
- **[매우 중요] 지엽적 변수명·기호 사용 금지 규칙**:
  - 해당 논문/특허 내부에서만 정의·사용되는 **지엽적(local) 변수명, 파라미터 기호, 제어변수 약어**를 요약 문장에 그대로 쓰지 말 것
  - 예를 들어, Dp, DH, DL, phi12, phi13, VoH, VoL, P2ref, P3ref, D1, D2, vT, LP, LS, LT, LM, n, k, F 같은 논문 내부 변수를 그대로 노출하면 독자가 이해할 수 없음
  - 이런 변수·기호가 등장하면, 그것이 **무엇을 의미하는지 개념적으로 한국어로 풀어서** 작성할 것
  - 잘못된 예: "제어변수 Dp, DH, DL, phi12, phi13 중 2개로 포트 전력 제약을 만족" → 변수명 나열이므로 금지
  - 올바른 예: "듀티비와 위상 변수 등 5개 제어변수 중 2개로 포트 전력 제약을 만족" → 변수를 개념적으로 설명
  - 잘못된 예: "입력 축을 VoH, VoL, P2ref, P3ref로 하여 4차원 룩업테이블 생성" → 변수명 나열이므로 금지
  - 올바른 예: "출력전압·포트전력 기준값을 입력 축으로 하는 4차원 룩업테이블 생성" → 변수를 개념적으로 설명
  - 잘못된 예: "D1,D2 관계와 VLink·VHV 변화에 무관한 vT 극성 기반 온오프 기준 적용" → 변수명 나열이므로 금지
  - 올바른 예: "듀티비 관계와 링크·고전압 버스 전압 변화에 무관한 변압기 극성 기반 온오프 기준 적용" → 변수를 개념적으로 설명
  - 단, 업계 표준 약어(MOSFET, PWM, PID, SOC, BMS, IGBT 등)나 학술 분야의 공인된 약어(GNSS, RTK, MIMO, FET 등)는 사용 가능
- 서술식 문단 금지, 개조식(짧은 문장/구)으로 작성. 개조식으로 작성하되, 내용이 구체적이고 상세해야 함
- 특허 문서의 경우: 부품·요소의 단순 참조 번호(예: 지지기35, 전축(38), 기판100 등)를 절대 포함하지 말 것. 부품·요소 이름만 쓰고 단순 참조 번호는 제거할 것. 단, 개선 효과의 수치(예: 16~18% 절감, 약 40% 향상)는 반드시 유지할 것.
- **중요** 요약 문장 끝에 "." 사용 금지
- 단위는 반드시 기호: 미터→m, 70도→70°, 퍼센트→%
- **[중요] 숫자와 단위 기호 사이에 띄어쓰기 금지**: 50 kHz→50kHz, 50 kW→50kW, 1.2 A→1.2A, 3.5 mm→3.5mm, 100 rpm→100rpm, 12 V→12V 등 모든 영문 단위는 숫자에 붙여 씀
- 불릿 접두 기호(•, -, ·, 번호)는 절대 사용하지 말고 내용만 출력
- 소수점은 반드시 18.4처럼 표기 ("점"으로 한글 사용 금지)

## 문서 유형 판별
- 논문이면 doc_type="paper", 특허이면 doc_type="patent"
- 특허 판별 키워드: claim, patent, 출원, 특허, 청구항, 발명

# 출력 형식
## 제목 (title) — '후보기술명' = 기술 개발 아이템 명칭으로 작성
- 이 연구/기술이 **어떤 새로운 시도와 컨셉**으로 **무엇을 개발·구현했는지**, 특징과 핵심이 잘 드러나는 **기술 개발 아이템 명칭**으로 작성
- 단순 논문 제목 번역이 아니라, 후보기술 개요서의 기술명으로서 한눈에 기술 핵심이 파악되도록 작성
- 목적 + 방법(기술)을 한 문장에 압축. 개선 효과/결과는 제목에 포함하지 말 것
- **[매우 중요] 제목은 반드시 기술 개발 아이템 명칭(기술/방법/시스템/장치/알고리즘/모델/구조/로직/센서/모듈/플랫폼 등)으로 끝나야 한다**
- **[금지] 제목이 개선 효과로 끝나면 안 됨. '~향상', '~개선', '~감소', '~절감', '~최적화', '~강화', '~확보', '~극복' 등 효과·결과를 나타내는 말로 끝내지 말 것**
- 공백 포함 50자 이내, 1문장, 개조식 (PPT 슬라이드 타이틀이므로 반드시 50자 이내로 작성)
- 아래 예시 스타일을 따라라 (모두 기술 아이템 명칭으로 끝남):
  • 경로 추종 지속여부 판정 로직 개발 → '개발'로 끝남 (28자)
  • 가변 속도 제안 협조 제어 방법 → '방법'으로 끝남 (16자)
  • ML 기반 NLOS 탐지·PPP-RTK 도심 정밀 측위 기술 → '기술'로 끝남 (23자)
- 잘못된 예시 (절대 금지):
  • ✗ "ML 기반 도심 측위 정확도 향상" → 효과('향상')로 끝남
  • ✗ "배선 길이 16~18% 절감" → 효과('절감')로 끝남
  • ✗ "NLOS 환경 위치 인식 성능 개선" → 효과('개선')로 끝남

## 헤드메시지 (head_messages)
- 정확히 2문장, 개조식. 명사형으로 문장 끝맺음하기. 문장 끝에 "." 사용 금지
- 각 문장은 **공백 포함 최대 200자 이내**로 간결하게 작성하여 PPT에서 문장바꿈이 발생하지 않도록 할 것
- 1번째 문장: 연구/기술 개발 목적 중심
- 2번째 문장: 해결방법 및 개선효과 중심
- 아래 요약된 문장 예시 스타일을 따라라:
  - 신호 특성 기반 머신러닝으로 NLOS와 거리 오차를 동시에 추정해 도심 환경에서도 스마트폰 PPP-RTK 차량 위치 정확도를 대폭 향상
  - ML 예측 결과를 확률적 관측 가중 모델로 변환해 기존 PPP-RTK 구조를 변경 없이 보강하는 실시간 차량 내비게이션 기술

## 기술목적 (purpose)
- **정확히 1-2개 불릿**, 개조식. 각 불릿은 1-2 문장 이내로 작성할 것
- 각 문장은 **공백 포함 최대 100자 이내**로 간결하게 작성하여 가급적 PPT에서 문장바꿈이 발생하지 않도록 할 것
- 아래 요약된 문장 예시 스타일을 따라라:
  - 도심 환경에서 NLOS와 다중경로로 저하되는 스마트폰 기반 PPP-RTK 차량 위치 인식 성능을 개선하고 추가 하드웨어 없이 실시간 적용 가능한차량 정밀 위치인식 기술을 구현
  - 도심 협곡(urban canyon) 환경에서 GNSS 신호 차단, 다중경로, NLOS로 인한 RTK/INS 위치 정확도 저하 문제 해결
  - 도심 협곡 환경에서 발생하는 GNSS NLOS·멀티패스 기반 연속적 이상 측정 문제 해결 

## 기존 기술의 문제점 (prior_problems)
- **정확히 1-2개 불릿**, 개조식. 각 불릿은 1~2문장 이내로 작성할 것
- 각 문장은 **공백 포함 최대 100자 이내**로 간결하게 작성하여 가급적 PPT에서 문장바꿈이 발생하지 않도록 할 것
- 아래 요약된 문장 예시 스타일을 따라라:
  - 기존 고정밀 측위(Multi-RTT 등)는 다수의 앵커(gNB)가 필수적이며, 단일 앵커 방식은 멀티패스 환경에서 각도 추정 오차로 성능이 저하됨
  - 기지국 신호 의존으로 인한 도심 협곡 및 터널 등 GNSS 음영 지역의 정확도 저하 
  - 단말 간 직접적인 위치 데이터 교환 체계 부재로 주변 차량의 동적 정보 반영의 한계
  - 기존 FGO 기반 GNSS/INS 연구는 느슨한 결합(loose coupling)에 머무르거나, 반송파 정수 모호수의 시간적 연속성을 충분히 활용하지 못함

## 제안 기술의 구체적인 방법 혹은 기술컨셉 (proposed_method) — 구체적이되, 개념 중심으로 작성
- 반드시 아래와 같은 JSON 배열(array) 형식으로 출력할 것 (dict/object 형식 절대 금지)
- 3~4개 상위 항목(title) × 각 3개 세부사항(details) 구조 (기본 ①~③개 상위 항목, 요약할 내용이 풍부하면 ①~④개 상위 항목까지 가능)
- 기존 기술 대비 무엇이 다른지, 핵심 메커니즘/구조/프로세스를 설명
- **주의**: 각 항목(title+details)은 개별 PPT 박스에 들어가므로, 각 세부사항(detail) 문장을 공백 포함 100자 이내로 작성할 것.
- **[매우 중요] 수식 절대 금지 규칙**:
  - 수식(수학식, 방정식, 수학 기호)을 그대로 나열하는 것을 **절대 금지**한다
  - |V0(1)|, sin(δ/2), Σwᵢyᵢ, x = f(y), Δx, ∫, ∂, argmax 등 수학적 표기·기호를 요약 문장에 포함하지 말 것
  - 수식이 있을 경우 해당 수식이 **의미하는 개선 포인트나 핵심 아이디어**를 일반인도 이해할 수 있는 **평이한 한국어**로 풀어서 설명할 것
  - 잘못된 예: "보상기로 |V0(1)|를 산출하고 sin(δ/2)로 매핑해 δ를 추정" → 수식 나열이므로 금지
  - 올바른 예: "전압 크기를 기반으로 위상 오차를 추정하는 보상 알고리즘 적용" → 수식 없이 개선 포인트를 설명
  - 잘못된 예: "x = Σwᵢyᵢ로 최종 값 산출" → 수식 나열이므로 금지
  - 올바른 예: "가중 합산 방식으로 최종 값을 산출" → 수식 없이 개념 설명
- 단순한 목적·효과가 아니라, 어떻게 구현되는지(프로세스), 어떤 핵심 기술 아이디어/발전 포인트가 적용되었는지를 중심으로 기술적으로 작성
- 알고리즘명·구조명·기법명 등 고유 명칭은 포함하되, 수식 자체를 복사하여 나열하지 않도록 주의
- **주의**: 배열의 `title` 값에 '1)', '2)' 와 같은 **번호나 기호를 절대 붙이지 마라**. (예: '반사체 탐지 및 위치 결정')
- proposed_method 형식 예시 (반드시 기호 없는 제목으로 작성):
  기본 3개 항목 예시 (대부분의 경우):
  [
    {"title": "① 상위 항목 1", "details": ["세부1", "세부2", "세부3"]},
    {"title": "② 상위 항목 2", "details": ["세부1", "세부2", "세부3"]},
    {"title": "③ 상위 항목 3", "details": ["세부1", "세부2", "세부3"]}
  ]
  내용이 충분히 풍부한 경우에만 4개 항목 (선택적):
  [
    {"title": "① 상위 항목 1", "details": ["세부1", "세부2", "세부3"]},
    {"title": "② 상위 항목 2", "details": ["세부1", "세부2", "세부3"]},
    {"title": "③ 상위 항목 3", "details": ["세부1", "세부2", "세부3"]},
    {"title": "④ 상위 항목 4", "details": ["세부1", "세부2", "세부3"]}
  ]
- 아래 예시 스타일을 따라라 (수식 없이, 개념·프로세스 중심):
  ① 반사체 탐지 및 위치 결정 (Sensing Phase)
  - RF 센싱 또는 LIDAR로 주변 건물 등 반사체의 위치를 사전에 파악하여 맵핑 데이터로 구축
  - 측정된 반사체 위치 정보를 위치 추정 엔티티(LMF 등)에 보고하여 이후 위치 추정의 기준점으로 활용
  ② 멀티패스 PRS(Positioning Reference Signal) 최적 구성
  - 반사체를 경유하는 경로를 포함하도록 PRS 전송 구성을 최적화하여, LOS뿐 아니라 NLOS 반사 경로까지 활용 가능하게 설계
  - 다각도 신호 도달 시간(ToA) 확보를 통해 단일 앵커 환경에서도 위치 추정 정밀도를 높이는 구조
  ③ 가상 앵커 기반 위치 추정 실행
  - 반사 경로 측정값을 별도 앵커에서 수신한 신호처럼 활용(Virtual Anchor 개념)하여 단일 노드에서도 삼변측량 기반 위치 추정 가능

## 개선효과 (improvements)
- **최대 3개 불릿**
- **[매우 중요]** 문서 내에 정량적 수치(%, dB, m, cm, 배수, 절감량 등)가 있으면 **반드시 해당 수치를 그대로 포함**하여 작성. 예: '배선 길이 16~18% 절감', '오차 40% 이상 감소'
- 수치가 전혀 없는 경우에만 정성적 설명으로 대체. 수치가 있는데 빠뜨리는 것은 절대 금지.
- 성능이 개선된 이유나 메커니즘 포함
- 각 불릿은 간결하게 핵심을 전달할 것 (1~2문장 이내)
- 아래 예시 스타일을 따라라:
  - 존 기반 IRN 재설계로 배선 길이·중량 16~18% 절감 및 열 방출 효율 향상
  - TK 단독 대비 ambiguity propagation 적용 시 수평 위치 정확도 약 40% 이상 개선
  - RTK/INS 융합 시 도시 협곡 환경에서도 수평 위치 오차 10–20 cm 수준 유지
  - GNSS 단절 구간에서 위치 RMSE: INS 단독 대비 최대 81% 감소, 가속도 기반 ML 대비 추가 39~57% 성능 개선

## 출처 정보 (paper_info) - 논문/특허 메타정보 
- journal_or_patent_office: 저널명 또는 특허청 (영어)
- paper_title: 논문/특허 제목 (영어)
- institution: 논문이면 연구기관(영어, 대학/회사명만), 특허이면 출원인(Assignee, 영어)을 기재
- doi_or_patent_no: 출처 고유 식별자
  - **논문**: 문서에 "DOI:"가 명시되어 있으면 DOI 값을 기재 (예: "10.1109/TVT.2024.1234567"). DOI가 없으면 논문 제목(영어)을 기재
  - **특허**: 특허번호를 **쉼표(,)나 공백 없이 붙여서** 기재 (예: "US1234567B1", "KR102345678B1", "WO2024123456A1"). 절대 "US 12,345,67 B1"처럼 구분하지 말 것
- year: 게재/출원 연도
- month: 게재/출원 월 (불명확하면 빈 문자열)

## 대표 이미지 (representative_figures)
- 제공된 PDF 문서에서 가장 대표적인 이미지 정확히 2개를 찾아 원본 그대로 추출해. 새로운 이미지를 생성하지 말고, 이미지 파일 형식으로 출력해.
- 첫번째 이미지 1개: 연구의 전체 컨셉/프레임워크 (role="concept")
- 두번째 이미지 1개: 연구의 효과/성능 (role="result")
- **[중요] 두 이미지는 반드시 서로 다른 Figure여야 한다. 같은 Figure를 두 번 선택하지 말 것.**
- **[중요] 동일 Figure의 단순 확대/반복/동일 그래프는 중복 선택 금지. 시각적 의미가 서로 다른 두 figure를 골라야 한다 (예: 구조도 1개 + 결과 그래프 1개).**
- 각각 page(페이지 번호), fig_number(Figure 번호), fig_title(Figure 제목) 기재
- fig_title은 논문 원문의 figure caption(Fig./Figure/그림 ~)을 그대로 사용. 원문 caption이 없으면 fig_title을 빈 문자열("")로 두고 임의로 만들지 말 것.

## Figure 캡션 (figure_captions_ko)
- 정확히 2개, 한글
- **캡션은 해당 그림이 무엇을 묘사하는 그림인지 설명하는 1문장**이다. 효과, 성능 향상, 개선 결과를 서술하는 문장이 아님. "~향상", "~개선", "~감소" 등 효과·결과를 나타내는 말로 캡션을 끝내지 말 것
- 원문의 Figure/Fig. 번호와 제목(caption)이 있으면 그것을 최우선으로 번역·활용하여 캡션 작성
- **[중요] 원문 caption을 찾지 못한 경우, 임의로 캡션을 지어내지 말고 빈 문자열("")로 둘 것. (이후 단계에서 fallback 문구로 자동 대체됨)**
- **[금지 표현]** 다음 표현은 캡션에 절대 사용하지 말 것: "본 발명은", "본 발명의", "본 발명에", "본 연구의", "본 연구에서", "본 논문의". 이런 표현 대신 기술 내용 자체를 직접 서술할 것
- **중요: 반드시 캡션은 명사로 끝내야 함**
- 원문 Figure 제목·설명을 최대한 반영하여 1문장, 30자 내외로 간결하게 작성.
- "그림1", "그림2", "Figure 1", "Figure 2"처럼 번호만 적힌 임시 라벨은 캡션으로 사용하지 말 것. 그림 내용 설명이 포함되어야 함.
- 예시: "MLP 기반 NLOS 탐지·거리 오차 예측과 PPP-RTK 결합 전체 프레임워크", "도심 환경별 PPP-RTK와 ML-보강 PPP-RTK 위치 오차 비교 결과표"

반드시 위 JSON 스키마에 맞는 JSON만 출력하라
"""


# ============================================================
# Text utilities
# ============================================================
def _normalize(text: str) -> str:
    t = (text or "").replace("\n", " ")
    return re.sub(r"\s{2,}", " ", t).strip()


def _strip_prefixes(text: str) -> str:
    t = _normalize(text)
    t = re.sub(r"^[•\-·\s]+", "", t).strip()
    t = re.sub(r"^\s*\d+\s*[\)\.]\s*", "", t).strip()
    return t


def _fix_punct(s: str) -> str:
    t = (s or "").strip()
    t = re.sub(r"(\d)\s*점\s*(\d)", r"\1.\2", t)
    t = re.sub(r"\s+점\s*$", "", t)
    return re.sub(r"\s{2,}", " ", t).strip()


def _fix_units(s: str) -> str:
    t = (s or "")
    # 한글 단위 → 기호 변환
    t = re.sub(r"\s*퍼센트\s*포인트\s*", "%p", t)
    t = re.sub(r"\s*퍼센트\s*", "%", t)
    t = re.sub(r"(\d)\s*도\b", r"\1°", t)
    t = re.sub(r"(\d)\s*킬로미터\b", r"\1km", t)
    t = re.sub(r"(\d)\s*미터\b", r"\1m", t)
    t = re.sub(r"(\d)\s*센티미터\b", r"\1cm", t)
    t = re.sub(r"(\d)\s*밀리미터\b", r"\1mm", t)
    # [영문 단위] 숫자와 단위 사이 공백 제거
    # 주파수/전력/전기
    t = re.sub(r"(\d)\s+(G?H?z|k?Hz|MHz|GHz|THz)", r"\1\2", t)
    t = re.sub(r"(\d)\s+(k?W|MW|GW|mW|μW|nW)", r"\1\2", t)
    t = re.sub(r"(\d)\s+(k?V|mV|μV)", r"\1\2", t)
    t = re.sub(r"(\d)\s+(k?A|mA|μA|nA)", r"\1\2", t)
    t = re.sub(r"(\d)\s+(k?Ω|MΩ|mΩ)", r"\1\2", t)
    t = re.sub(r"(\d)\s+(dBm?|dBi|dBc|dBA?)", r"\1\2", t)
    t = re.sub(r"(\d)\s+(k?Wh|mAh|Ah)", r"\1\2", t)
    # 길이
    t = re.sub(r"(\d)\s+(km|cm|mm|nm|μm|pm)", r"\1\2", t)
    # 속도/압력/온도
    t = re.sub(r"(\d)\s+rpm\b", r"\1rpm", t)
    t = re.sub(r"(\d)\s+(k?m/s|km/h)", r"\1\2", t)
    t = re.sub(r"(\d)\s+(MPa|GPa|kPa|hPa|Pa)", r"\1\2", t)
    t = re.sub(r"(\d)\s+°C\b", r"\1°C", t)
    # ppm / bit
    t = re.sub(r"(\d)\s+ppm\b", r"\1ppm", t)
    t = re.sub(r"(\d)\s+(Gbps|Mbps|kbps|bps)", r"\1\2", t)
    return re.sub(r"\s{2,}", " ", t).strip()


def clean(s: str, max_chars: int = 95, hard_cut: bool = False) -> str:
    """기본 텍스트 정리. hard_cut=True이면 강제 절단(내부용).
    일반적으로는 10자 이상 초과할 때만 절단 대상으로 표시하고,
    실제 API 재요약은 normalize_brief에서 수행."""
    t = _fix_units(_fix_punct(_strip_prefixes(s)))
    if len(t) <= max_chars:
        return t
    if hard_cut:
        cut = t[:max_chars]
        if " " in cut:
            cut2 = cut.rsplit(" ", 1)[0].rstrip()
            if len(cut2) >= max_chars - 15:
                return cut2
        return cut.rstrip()
    # 10자 미만 초과: 그냥 통과 (사람이 검수)
    if len(t) <= max_chars + 10:
        return t
    # 10자 이상 초과: 원문 보존, 재요약 대상으로 표시 (호출자에서 처리)
    return t


def _ensure_list(x, min_n=0, max_n=99, placeholder="") -> list:
    arr = x if isinstance(x, list) else []
    arr = [clean(s) for s in arr if clean(s)]
    arr = arr[:max_n]
    while len(arr) < min_n:
        arr.append(placeholder)
    return arr


def shorten_text_via_llm(text: str, max_chars: int, model: str = DEFAULT_MODEL) -> str:
    """글자 수 제한을 초과한 문장을 API로 재요약. 실패 시 강제 절단 반환.
    
    NOTE: 재요약 전용이라 reasoning 모델(gpt-5, o1, o3)은 부적합 → 항상 gpt-4o-mini 사용.
    Reasoning 모델은 max_completion_tokens 중 상당량을 내부 사고에 소비해 빈 응답을 반환할 수 있음.
    """
    api_key = os.getenv("OPENAI_API_KEY", "").strip()
    gemini_key = os.getenv("GEMINI_API_KEY", "").strip()

    # 재요약은 무조건 빠르고 안정적인 모델 사용 (reasoning 모델 제외)
    shorten_model = model
    if IS_GEMINI_MODEL(model):
        # Gemini면 flash로 변경 (빠름)
        shorten_model = "gemini-2.0-flash"
    elif any(x in model.lower() for x in ["o1", "o3", "gpt-5", "gpt-o3"]):
        # Reasoning 모델이면 gpt-4o-mini로 다운그레이드
        shorten_model = "gpt-4o-mini"

    prompt = (
        f"다음 문장을 {max_chars}자 이내로 줄여줘. 의미는 최대한 유지하고, 개조식으로 작성해. "
        f"불릿 기호(•, -, ·) 없이 내용만 출력해. 문장 끝에 마침표(.) 사용 금지.\n\n"
        f"원문: {text}"
    )
    try:
        if IS_GEMINI_MODEL(shorten_model) and HAS_GEMINI and gemini_key:
            genai.configure(api_key=gemini_key)
            gmodel = genai.GenerativeModel(shorten_model)
            resp = gmodel.generate_content(prompt)
            result = (resp.text or "").strip()
        else:
            if not api_key:
                raise RuntimeError("OPENAI_API_KEY not set")
            client = OpenAI(api_key=api_key)
            params: Dict[str, Any] = {
                "model": shorten_model,
                "messages": [{"role": "user", "content": prompt}],
                "max_completion_tokens": 500,
                "temperature": 0.3,
            }
            resp = client.chat.completions.create(**params)
            result = (resp.choices[0].message.content or "").strip()

        # 빈 응답 가드: 빈 결과면 즉시 hard-cut 반환
        if not result:
            print(f"  [WARN] shorten_text_via_llm returned empty. Hard cutting.")
            return clean(text, max_chars, hard_cut=True)

        result = _fix_units(_fix_punct(_strip_prefixes(result)))

        # strip 후에도 비어있으면 hard-cut
        if not result:
            return clean(text, max_chars, hard_cut=True)

        # 여전히 초과하면 강제 절단
        if len(result) > max_chars + 10:
            return clean(result, max_chars, hard_cut=True)
        return result
    except Exception as e:
        print(f"  [WARN] shorten_text_via_llm failed: {e}. Hard cutting.")
        return clean(text, max_chars, hard_cut=True)


def _strip_patent_object_numbers(text: str) -> str:
    """특허 본문의 구성 요소 참조 번호만 제거. 수치/단위(%, m, %, 절감량 등)는 보존.
    - 괄호형: 지지기(35), 전축(38, 39) → 지지기, 전축  [단, (16~18%)처럼 단위가 포함되면 유지]
    - 직접형: 지지기35, 전축38 → 지지기, 전축
    - 혼합형: 기판100a → 기판
    - 수치 표현(16~18%, 약 40%, 10 cm)은 절대 제거하지 않음
    """
    # 1) 괄호형: "객체명(번호)" — 단, 퍼센트(%), 단위(m, cm, mm, dB, Hz 등)가 포함된 경우 보존
    #    예: 지지기(35) → 지지기 | 절감(16~18%) → 절감(16~18%) 보존
    t = re.sub(
        r'\s*\(\s*(\d+[a-zA-Z]?(?:\s*,\s*\d+[a-zA-Z]?)*)\s*\)',
        lambda m: '' if not re.search(r'[%~\-]|[a-zA-Z]{2,}', m.group(1)) else m.group(0),
        text
    )
    # 2) 직접형: 한글 바로 뒤에 붙은 단순 숫자 제거 — "지지기35" → "지지기"
    #    단, 퍼센트 기호나 단위가 바로 뒤에 오는 경우(예: "절감16%") 는 건너뜀
    t = re.sub(r'(?<=[가-힣])(\d+[a-zA-Z]?)(?!\s*[%㎝㎞㎜°]|\s*(?:cm|mm|km|dB|Hz|kHz|MHz|GHz|%|px))', '', t)
    t = re.sub(r'\s{2,}', ' ', t).strip()
    return t


def pick_primary_affiliation(aff: str) -> str:
    s = (aff or "").strip()
    if not s:
        return ""
    first = re.split(r"\s*(?:;|\|/|\band\b|&)\s*", s, maxsplit=1, flags=re.IGNORECASE)[0].strip()
    parts = [p.strip() for p in first.split(",") if p.strip()]
    exclude = re.compile(
        r"\b(department|dept\.?|school|faculty|lab(oratory)?|center|centre|division|"
        r"unit|group|program|college|institute|research\s*center|graduate\s*school)\b",
        re.IGNORECASE,
    )
    kept = [p for p in parts if not exclude.search(p)]
    geo = re.compile(
        r"\b(republic of korea|south korea|korea|usa|united states|china|japan|germany|"
        r"france|uk|england|seoul|busan|tokyo|beijing|berlin|paris|london|singapore)\b",
        re.IGNORECASE,
    )
    kept2 = [p for p in kept if not geo.search(p)]
    if kept2:
        return kept2[-1].strip()
    if kept:
        return kept[-1].strip()
    if parts:
        return parts[-1].strip()
    return first


# ============================================================
# PDF extraction
# ============================================================
def extract_text(doc: fitz.Document, max_pages: int = 20) -> str:
    pages = min(max_pages, doc.page_count)
    chunks = []
    for i in range(pages):
        try:
            chunks.append(doc.load_page(i).get_text("text"))
        except Exception:
            continue
    text = "\n".join(chunks)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    return text


# ============================================================
# Image extraction (from v1.2)
# ============================================================
def trim_white_margin(png_bytes: bytes, pad: int = 6) -> bytes:
    if Image is None:
        return png_bytes
    try:
        im = Image.open(io.BytesIO(png_bytes)).convert("RGB")
    except Exception:
        return png_bytes
    w, h = im.size
    px = im.load()
    step = max(1, min(w, h) // 350)
    thr = 245

    def row_ink(y):
        for x in range(0, w, step):
            r, g, b = px[x, y]
            if r < thr or g < thr or b < thr:
                return True
        return False

    def col_ink(x, y0, y1):
        for y in range(y0, y1, step):
            r, g, b = px[x, y]
            if r < thr or g < thr or b < thr:
                return True
        return False

    top = 0
    for y in range(0, h, step):
        if row_ink(y):
            top = y; break
    else:
        return png_bytes

    bottom = h - 1
    for y in range(h - 1, -1, -step):
        if row_ink(y):
            bottom = y; break

    left = 0
    for x in range(0, w, step):
        if col_ink(x, top, bottom + 1):
            left = x; break

    right = w - 1
    for x in range(w - 1, -1, -step):
        if col_ink(x, top, bottom + 1):
            right = x; break

    left = max(0, left - pad)
    top = max(0, top - pad)
    right = min(w - 1, right + pad)
    bottom = min(h - 1, bottom + pad)

    if (right - left) < int(w * 0.15) or (bottom - top) < int(h * 0.15):
        return png_bytes

    cropped = im.crop((left, top, right + 1, bottom + 1))
    out = io.BytesIO()
    cropped.save(out, format="PNG", optimize=True)
    return out.getvalue()


_CAPTION_START_RE = re.compile(
    r"^\s*(Fig\.?|Figure|그림|도면)\s*(\d+)([a-zA-Z])?\s*[\.\:\)\-\—\–]\s*",
    re.IGNORECASE,
)
_SUBLABEL_ONLY_RE = re.compile(r"^\s*\(?([a-zA-Z])\)?\s*$")
# 본문에서 figure를 참조하는 문장 (캡션이 아님). 'Fig. N' 뒤에 동사/접속사/관계사 등이 오면 본문 참조.
# 예: "Fig. 1 shows ...", "Fig. 8b depicts ...", "Fig. 13b where the ...", "Fig. 9. Clearly, from ..."
_BODY_REF_AFTER_NUM_RE = re.compile(
    r"^\s*(?:Fig\.?|Figure|그림|도면)\s*\d+[a-zA-Z]?\s*[\.\,]?\s*"
    r"(?:shows?|depicts?|displays?|presents?|illustrates?|describes?|indicates?|reveals?|"
    r"demonstrates?|provides?|gives?|is|are|was|were|will|would|can|could|may|might|"
    r"clearly|where|when|while|after|before|above|below|here|there|in|on|of|for|"
    r"and|but|so|because|since|although|however)\b",
    re.IGNORECASE,
)


def _find_figure_captions_for_page(page) -> List[Dict]:
    """페이지에서 Figure 캡션 텍스트와 위치(bbox)를 추출하여 반환.
    반환: [{"text": str, "bbox": [x0,y0,x1,y1], "y": float, "figure_number": "Fig. 5"}, ...]

    개선 사항 (요청서 §5):
      1) 라인 단위로 스캔하여 Fig./Figure/그림/도면 + 숫자 패턴을 정규식으로 탐지
      2) 캡션 다음에 이어지는 줄을 자동 병합(다음 캡션/섹션/너무 먼 텍스트 전까지)
      3) 단독 (a)/(b) 라벨은 캡션 시작으로 인정하지 않음 — 실제 Fig. N. 본문을 캡션으로 사용
      4) caption_bbox 와 figure_number(예: 'Fig. 5') 를 함께 반환
    """
    captions: List[Dict] = []
    try:
        blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE).get("blocks", [])
    except Exception:
        return captions

    # 1) 모든 텍스트 라인을 평탄화 (bbox 포함)
    lines: List[Dict] = []
    for block in blocks:
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            line_text = "".join(span.get("text", "") for span in line.get("spans", [])).strip()
            if not line_text:
                continue
            lbbox = line.get("bbox")
            if not lbbox:
                continue
            lines.append({"text": line_text, "bbox": list(lbbox)})

    if not lines:
        return captions

    # top → bottom, left → right 순으로 정렬
    lines.sort(key=lambda l: (round(l["bbox"][1], 1), l["bbox"][0]))

    # 2) 캡션 시작 라인을 찾고, 이어지는 줄을 병합
    n = len(lines)
    i = 0
    while i < n:
        ln = lines[i]
        text = ln["text"]
        m = _CAPTION_START_RE.match(text)
        if not m:
            i += 1
            continue
        # 본문 텍스트에서 figure를 참조하는 문장은 캡션으로 인정하지 않는다.
        # 예: "Fig. 1 shows ...", "Fig. 8b depicts ...", "Fig. 13b where ...".
        if _BODY_REF_AFTER_NUM_RE.match(text):
            i += 1
            continue
        prefix_raw = m.group(1).lower()
        if prefix_raw.startswith("fig"):
            fig_label = "Fig."
        elif prefix_raw == "figure":
            fig_label = "Figure"
        elif prefix_raw == "그림":
            fig_label = "그림"
        elif prefix_raw == "도면":
            fig_label = "도면"
        else:
            fig_label = m.group(1)
        fig_num = m.group(2)
        sub = (m.group(3) or "").strip()
        figure_number = f"{fig_label} {fig_num}{sub}".strip()

        cap_bbox = list(ln["bbox"])
        cap_text_parts = [text]

        j = i + 1
        line_height = max(cap_bbox[3] - cap_bbox[1], 1.0)
        while j < n and len(cap_text_parts) < 6:
            nxt = lines[j]
            ntxt = nxt["text"]
            # 새 캡션 패턴이 나오면 중단
            if _CAPTION_START_RE.match(ntxt):
                break
            # (a) / (b) 등 단독 sublabel은 캡션 일부로 끌어오지 않음
            if _SUBLABEL_ONLY_RE.match(ntxt):
                break
            # 수직 거리 검사 — 다음 줄 top이 현재 캡션 bottom으로부터 line_height*1.6 이내
            gap = nxt["bbox"][1] - cap_bbox[3]
            if gap > line_height * 1.6:
                break
            # 수평 겹침 검사 — 캡션 x-범위와 어느 정도 겹쳐야 함
            cx0, cx1 = cap_bbox[0], cap_bbox[2]
            nx0, nx1 = nxt["bbox"][0], nxt["bbox"][2]
            if min(cx1, nx1) - max(cx0, nx0) < -10:
                break
            # 섹션 헤더처럼 보이면 중단(전부 대문자, 짧음)
            if len(ntxt) < 30 and re.match(r"^[A-Z][A-Z0-9 \.\-]+$", ntxt):
                break
            cap_text_parts.append(ntxt)
            cap_bbox[0] = min(cap_bbox[0], nxt["bbox"][0])
            cap_bbox[1] = min(cap_bbox[1], nxt["bbox"][1])
            cap_bbox[2] = max(cap_bbox[2], nxt["bbox"][2])
            cap_bbox[3] = max(cap_bbox[3], nxt["bbox"][3])
            line_height = max(nxt["bbox"][3] - nxt["bbox"][1], line_height)
            j += 1

        merged_text = " ".join(cap_text_parts)
        merged_text = re.sub(r"\s+", " ", merged_text).strip()
        captions.append({
            "text": merged_text,
            "bbox": cap_bbox,
            "y": cap_bbox[1],
            "figure_number": figure_number,
        })
        i = j

    return captions


def _compute_overlap_ratio(boxA, boxB) -> float:
    """두 BBox 간의 겹치는 영역 비율(IoU 형태)을 계산합니다."""
    xA = max(boxA[0], boxB[0])
    yA = max(boxA[1], boxB[1])
    xB = min(boxA[2], boxB[2])
    yB = min(boxA[3], boxB[3])

    interArea = max(0, xB - xA) * max(0, yB - yA)
    if interArea == 0:
        return 0.0
    boxAArea = (boxA[2] - boxA[0]) * (boxA[3] - boxA[1])
    boxBArea = (boxB[2] - boxB[0]) * (boxB[3] - boxB[1])
    
    # 두 박스 중 더 작은 박스 대비 겹치는 비율을 반환 (포함되는 경우 1.0)
    min_area = min(boxAArea, boxBArea)
    if min_area == 0:
        return 0.0
    return interArea / float(min_area)


def _extract_surrounding_text(page, img_bbox, margin=150) -> str:
    """이미지 Bbox 주변(상/하단)의 문맥 텍스트를 추출합니다."""
    if not img_bbox:
        return ""
    surrounding = []
    try:
        blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE).get("blocks", [])
        for block in blocks:
            if block.get("type") != 0:
                continue
            b_bbox = block.get("bbox")
            if not b_bbox: continue
            
            # y축 거리 계산: 이미지 위(margin 이내) 또는 아래(margin 이내)
            dist_top = abs(img_bbox[1] - b_bbox[3])
            dist_bottom = abs(img_bbox[3] - b_bbox[1])
            dist_center = abs((img_bbox[1]+img_bbox[3])/2 - (b_bbox[1]+b_bbox[3])/2)
            
            if dist_top <= margin or dist_bottom <= margin or dist_center <= margin + (img_bbox[3]-img_bbox[1])/2:
                block_text = "".join(
                    span["text"] for line in block.get("lines", []) for span in line.get("spans", [])
                ).strip()
                if block_text and len(block_text) > 5:
                    surrounding.append(block_text)
    except Exception:
        pass
    return "\n".join(surrounding)


def extract_images(doc: fitz.Document, max_candidates=20, min_pixels=120_000) -> List[Dict]:
    import hashlib
    candidates = []
    seen_xrefs = set()
    seen_hashes = set()  # 이미지 바이트 해시 기반 중복 제거
    seen_bboxes_per_page = {}  # 페이지별로 이미 추출된 이미지의 BBox 기록 (겹침 방지)
    seen_size_byte_keys = set()  # (width, height, byte_len) 동일 이미지 중복 차단
    figure_index_counter = 0  # 추출 순서 기반 전역 figure_index
    for pi in range(doc.page_count):
        page = doc.load_page(pi)
        try:
            page_text = page.get_text("text") or ""
        except Exception:
            page_text = ""

        # 페이지 전체 캡션 힌트 (기존 방식 — 폴백용)
        cap_lines = [ln.strip() for ln in page_text.splitlines()
                     if ln.strip() and any(ln.strip().lower().startswith(p) for p in
                                           ["fig.", "figure", "그림", "표", "도면"])]
        page_caption_hint = " | ".join(cap_lines[:6])

        # 위치 기반 캡션 목록 추출
        page_captions = _find_figure_captions_for_page(page)

        try:
            imgs = page.get_images(full=True)
        except Exception:
            imgs = []

        # 각 이미지의 페이지 내 위치(bbox) 파악
        img_bboxes = {}
        try:
            for item in page.get_image_info(xrefs=True):
                xref_val = item.get("xref", 0)
                bbox = item.get("bbox")  # (x0, y0, x1, y1)
                if bbox and xref_val:
                    img_bboxes[xref_val] = bbox
        except Exception:
            pass

        for img in imgs:
            xref = img[0]
            if xref in seen_xrefs:
                continue
            seen_xrefs.add(xref)
            try:
                pix = fitz.Pixmap(doc, xref)
                if pix.n >= 5:
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                area = pix.width * pix.height
                if area < min_pixels:
                    continue
                png_bytes = trim_white_margin(pix.tobytes("png"))
                # 트리밍 후 실제 width/height — PIL 사용 가능 시 우선 사용
                img_w, img_h = pix.width, pix.height
                if Image:
                    im = Image.open(io.BytesIO(png_bytes))
                    img_w, img_h = im.size[0], im.size[1]
                    area = img_w * img_h

                # 이미지 바이트 해시로 중복 검사 (다른 xref지만 동일 이미지 제거)
                img_hash = hashlib.md5(png_bytes).hexdigest()
                if img_hash in seen_hashes:
                    continue

                # 동일 (width, height, byte_len) 조합 → 거의 확실한 동일 이미지
                size_byte_key = (img_w, img_h, len(png_bytes))
                if size_byte_key in seen_size_byte_keys:
                    continue

                img_bbox = img_bboxes.get(xref)

                # BBox 기반 중복 검사: 같은 페이지 내에서 이미 추출한 이미지 영역과 크게 겹치면 무시
                is_overlap = False
                if img_bbox:
                    if pi not in seen_bboxes_per_page:
                        seen_bboxes_per_page[pi] = []
                    for existing_bbox in seen_bboxes_per_page[pi]:
                        if _compute_overlap_ratio(img_bbox, existing_bbox) > 0.8:
                            is_overlap = True
                            break
                    if is_overlap:
                        continue
                    seen_bboxes_per_page[pi].append(img_bbox)

                seen_hashes.add(img_hash)
                seen_size_byte_keys.add(size_byte_key)

                # 이미지 위치 기반으로 가장 가까운 캡션 및 주변 문맥 찾기
                # (요청서 §5: 이미지 bbox 바로 아래 텍스트 우선, x축 겹침 고려, 너무 먼 텍스트 제외)
                nearest_caption = ""
                caption_source = ""
                caption_bbox = None
                figure_number_str = ""
                surrounding_text = ""

                if img_bbox and page_captions:
                    ix0, iy0, ix1, iy1 = img_bbox
                    img_w_span = max(ix1 - ix0, 1.0)

                    def _x_overlap_ratio(cb):
                        if not cb:
                            return 0.0
                        left = max(ix0, cb[0])
                        right = min(ix1, cb[2])
                        return max(0.0, right - left) / img_w_span

                    # 1순위: 이미지 바로 아래 (cap.y_top > img.y_bottom 근처) + x-겹침
                    below_candidates = []
                    for c in page_captions:
                        cb = c.get("bbox")
                        cy_top = cb[1] if cb else c["y"]
                        if cy_top >= iy1 - 5:  # 이미지 하단보다 아래(약간의 여유)
                            dist = cy_top - iy1
                            xov = _x_overlap_ratio(cb)
                            below_candidates.append((dist, -xov, c))
                    if below_candidates:
                        # 너무 먼 caption(>250pt)은 제외
                        below_candidates = [b for b in below_candidates if b[0] <= 250]
                    if below_candidates:
                        # x-겹침 큰 것 우선, 그 다음 가까운 것
                        below_candidates.sort(key=lambda t: (t[1], t[0]))
                        chosen = below_candidates[0][2]
                        nearest_caption = chosen["text"]
                        caption_bbox = chosen.get("bbox")
                        figure_number_str = chosen.get("figure_number", "")
                        caption_source = "below_image_text"
                    else:
                        # 2순위: 이미지 위쪽 (cap.y_bottom < img.y_top 근처)
                        above_candidates = []
                        for c in page_captions:
                            cb = c.get("bbox")
                            cy_bottom = cb[3] if cb else c["y"]
                            if cy_bottom <= iy0 + 5:
                                dist = iy0 - cy_bottom
                                xov = _x_overlap_ratio(cb)
                                above_candidates.append((dist, -xov, c))
                        above_candidates = [a for a in above_candidates if a[0] <= 250]
                        if above_candidates:
                            above_candidates.sort(key=lambda t: (t[1], t[0]))
                            chosen = above_candidates[0][2]
                            nearest_caption = chosen["text"]
                            caption_bbox = chosen.get("bbox")
                            figure_number_str = chosen.get("figure_number", "")
                            caption_source = "above_image_text"

                if img_bbox:
                    surrounding_text = _extract_surrounding_text(page, img_bbox)

                # caption_hint는 Vision LLM 컨텍스트용 — nearest_caption이 있으면 그것만 사용
                # (page-level dump 사용을 줄여 다른 figure caption과 섞이지 않도록)
                if nearest_caption:
                    caption_hint = nearest_caption
                else:
                    caption_hint = page_caption_hint
                    if page_caption_hint and not caption_source:
                        caption_source = "page_caption_hint"

                # figure_number_str(예: 'Fig. 5') 가 있으면 우선, 없으면 캡션 텍스트에서 파싱
                fig_label = (_normalize_fig_number(figure_number_str)
                             or _parse_fig_number_from_caption(nearest_caption)
                             or "")

                # 요청서 §4,§6 — category 분류 및 figure_number 검증
                _cat, _role, _score = _classify_figure_category(nearest_caption)
                # image_figure_number: 이미지 bbox 매칭에서 얻은 figure_number
                _img_fig_num = figure_number_str or (f"Fig. {fig_label}" if fig_label else "")
                # caption_figure_number: 캡션 텍스트에서 파싱한 figure_number
                _cap_fig_num_raw = _parse_fig_number_from_caption(nearest_caption)
                _cap_fig_num = f"Fig. {_cap_fig_num_raw}" if _cap_fig_num_raw else ""

                # §6: image_figure_number와 caption_figure_number 불일치 시 경고
                _fig_mismatch = False
                if _img_fig_num and _cap_fig_num:
                    if _normalize_fig_number(_img_fig_num) != _normalize_fig_number(_cap_fig_num):
                        _fig_mismatch = True
                        print(f"  [WARN] figure-caption mismatch: image={_img_fig_num}, "
                              f"caption={_cap_fig_num}, page={pi+1}. Rejecting caption.")
                        # 매칭 오류 시 캡션을 제거하고 figure_number만 유지
                        nearest_caption = ""
                        caption_hint = page_caption_hint or ""
                        caption_source = "mismatch_cleared"

                figure_index_counter += 1
                candidates.append({
                    "page": pi + 1,
                    "page_number": pi + 1,           # alias — 외부 명세 호환
                    "figure_index": figure_index_counter,
                    "area": area,
                    "width": img_w,
                    "height": img_h,
                    "bbox": list(img_bbox) if img_bbox else None,
                    "group_bbox": list(img_bbox) if img_bbox else None,
                    "png_bytes": png_bytes,
                    "byte_len": len(png_bytes),
                    "caption_hint": caption_hint,
                    "nearest_caption": nearest_caption,
                    # 요청서 §4 — 명시적 키 (외부/검증용)
                    "original_caption": nearest_caption,
                    "translated_caption": "",  # 번역은 후속 단계에서 채움
                    "caption_bbox": caption_bbox,
                    "caption_source": caption_source,  # below_image_text / above_image_text / page_caption_hint / ""
                    "crop_source": "embedded_image",
                    "contains_real_figure": True,    # embedded image는 실제 raster image
                    "contains_body_text": False,
                    "is_text_only": False,
                    "surrounding_text": surrounding_text,
                    "fig_label": fig_label,           # 순수 숫자 (예: "3", "1a")
                    "figure_number": _img_fig_num,
                    # 요청서 §6 — 매칭 검증용
                    "image_figure_number": _img_fig_num,
                    "caption_figure_number": _cap_fig_num,
                    # 요청서 §4 — category 분류
                    "category": _cat,
                    "role_candidate": _role,
                    "selection_score": round(_score, 3),
                    "hash": img_hash,                 # 외부 명세 호환
                    "_img_hash": img_hash,            # 기존 코드 호환 (제거하지 않음)
                    "image_hash": img_hash,           # 요청서 §4 명시 키
                    "image_source": "embedded",
                    "selection_reason": "",
                })
            except Exception:
                continue

    candidates.sort(key=lambda d: d["area"], reverse=True)

    # ── fig_label 기준 중복 제거 ──────────────────────────────────────────
    # 같은 Figure가 여러 sub-image로 분해된 경우, area가 가장 큰 것 하나만 남김
    # → max_candidates 슬롯을 특정 Figure의 파편들이 독점하는 문제 방지
    seen_fig_labels: set = set()
    deduped: List[Dict] = []
    no_label: List[Dict] = []  # fig_label 없는 이미지 (별도 처리)
    for c in candidates:
        lbl = c.get("fig_label", "")
        if lbl:
            if lbl not in seen_fig_labels:
                seen_fig_labels.add(lbl)
                deduped.append(c)
            # else: 같은 fig_label의 더 작은 이미지 → 스킵
        else:
            no_label.append(c)  # 캡션 없는 이미지는 나중에 빈 슬롯 채우기

    # 캡션 있는 Figure 먼저, 이후 캡션 없는 것 추가
    merged = deduped + no_label

    # ── 요청서 §6: page render + caption-bbox 기준 crop fallback ────────────
    # embedded image 추출만으로는 못 잡은 figure(벡터 그래픽 등)를 복구한다.
    # 이미 캡션과 매칭된 figure_number는 건너뛰고, 미매칭 캡션에 대해서만 crop.
    try:
        crop_candidates = _extract_figures_via_page_crop(doc, merged)
    except Exception as e:
        print(f"  [WARN] page-crop fallback failed entirely: {e}")
        crop_candidates = []

    if crop_candidates:
        existing_hashes = {c.get("hash") or c.get("_img_hash") for c in merged if c.get("hash") or c.get("_img_hash")}
        existing_fignum = {(_normalize_fig_number(c.get("figure_number","") or c.get("fig_label","")) or "")
                           for c in merged}
        existing_fignum.discard("")
        for cc in crop_candidates:
            ch = cc.get("hash") or cc.get("_img_hash")
            cf = _normalize_fig_number(cc.get("figure_number","") or cc.get("fig_label","")) or ""
            if ch and ch in existing_hashes:
                continue
            if cf and cf in existing_fignum:
                # 같은 figure_number는 embedded image를 이미 보유 → crop 본은 폐기
                continue
            figure_index_counter += 1
            cc["figure_index"] = figure_index_counter
            merged.append(cc)
            if ch:
                existing_hashes.add(ch)
            if cf:
                existing_fignum.add(cf)

    # 요청서 §10 — text-only 후보 제거 (contains_real_figure=False)
    rejected = [c for c in merged if not c.get("contains_real_figure", True)]
    for r in rejected:
        print(f"  [WARN] text-only candidate rejected: page={r.get('page','?')}, "
              f"fig={r.get('figure_number','?')}, "
              f"crop_source={r.get('crop_source','?')}")
    merged = [c for c in merged if c.get("contains_real_figure", True)]

    final = merged[:max_candidates]

    # 요청서 §13 — figure candidate 로그
    print(f"  [INFO] figure candidates: total={len(final)} "
          f"(embedded={sum(1 for c in final if c.get('image_source')=='embedded')}, "
          f"page_crop={sum(1 for c in final if c.get('image_source')=='page_crop')})")
    for c in final:
        figno = c.get("figure_number") or (f"Fig. {c.get('fig_label')}" if c.get('fig_label') else "?")
        oc = (c.get("original_caption") or "")[:80]
        src = c.get("caption_source") or "?"
        idx_str = str(c.get("figure_index") if c.get("figure_index") is not None else "?")
        page_str = str(c.get("page") if c.get("page") is not None else "?")
        print(f"  [INFO]   - idx={idx_str:>2} page={page_str:>2} "
              f"src={c.get('image_source','?')} {figno!s:<10} caption_src={src} "
              f"original_caption='{oc}'")

    return final


def _x_overlap(a: List[float], b: List[float]) -> float:
    """두 bbox의 x축 겹침 길이(pt). 양수면 겹침, 음수면 분리."""
    return min(a[2], b[2]) - max(a[0], b[0])


def _collect_visual_rects(page, region_bbox: List[float]) -> List[List[float]]:
    """페이지 region 내부의 시각 객체(이미지·드로잉) bbox 목록.
    text block은 제외 — figure visual content만 모은다.
    region 내부와 일정 비율 이상 겹치는 것만 반환.
    """
    rects: List[List[float]] = []
    rx0, ry0, rx1, ry1 = region_bbox
    region_area = max((rx1 - rx0) * (ry1 - ry0), 1.0)

    # 1) 이미지 객체
    try:
        for it in page.get_image_info(xrefs=True):
            bb = it.get("bbox")
            if not bb:
                continue
            ix = max(0.0, min(rx1, bb[2]) - max(rx0, bb[0]))
            iy = max(0.0, min(ry1, bb[3]) - max(ry0, bb[1]))
            if ix * iy > 0.0:
                rects.append([bb[0], bb[1], bb[2], bb[3]])
    except Exception:
        pass

    # 2) text dict의 image block (type=1)
    try:
        blocks = page.get_text("dict").get("blocks", [])
        for blk in blocks:
            if blk.get("type") != 1:
                continue
            bb = blk.get("bbox")
            if not bb:
                continue
            ix = max(0.0, min(rx1, bb[2]) - max(rx0, bb[0]))
            iy = max(0.0, min(ry1, bb[3]) - max(ry0, bb[1]))
            if ix * iy > 0.0:
                rects.append([bb[0], bb[1], bb[2], bb[3]])
    except Exception:
        pass

    # 3) drawings (벡터 그래픽)
    try:
        for d in page.get_drawings():
            r = d.get("rect")
            if not r:
                continue
            ix = max(0.0, min(rx1, r.x1) - max(rx0, r.x0))
            iy = max(0.0, min(ry1, r.y1) - max(ry0, r.y0))
            if ix * iy > 0.0:
                rects.append([r.x0, r.y0, r.x1, r.y1])
    except Exception:
        pass

    return rects


def _collect_text_rects(page, region_bbox: List[float], min_chars: int = 8) -> List[List[float]]:
    """페이지 region 내부의 본문 텍스트 블록 bbox 목록 (figure 내부의 짧은 라벨은 제외)."""
    rects: List[List[float]] = []
    rx0, ry0, rx1, ry1 = region_bbox
    try:
        blocks = page.get_text("dict").get("blocks", [])
        for blk in blocks:
            if blk.get("type") != 0:
                continue
            text = "".join(
                span.get("text", "") for line in blk.get("lines", [])
                for span in line.get("spans", [])
            ).strip()
            if len(text) < min_chars:
                continue
            bb = blk.get("bbox")
            if not bb:
                continue
            cx = (bb[0] + bb[2]) / 2.0
            cy = (bb[1] + bb[3]) / 2.0
            if rx0 <= cx <= rx1 and ry0 <= cy <= ry1:
                rects.append([bb[0], bb[1], bb[2], bb[3]])
    except Exception:
        pass
    return rects


def _compute_figure_visual_bbox(page, cap_bbox: List[float], upper_limit: float,
                                 column_left: float, column_right: float) -> Optional[List[float]]:
    """caption_bbox 바로 위쪽 figure visual content의 정확한 bbox 계산.

    1) caption 위쪽 column 영역 안에서 시각 객체(이미지·드로잉)를 모두 모은다.
    2) caption에 가장 가까운(아래쪽) 클러스터를 figure 본체로 본다.
    3) 그 클러스터의 union bbox를 반환. 시각 객체가 없으면 None.

    Returns: [x0, y0, x1, y1] or None
    """
    # 검색 region: caption 위쪽 column 폭 전체
    region = [column_left, max(0.0, upper_limit), column_right, max(0.0, cap_bbox[1] - 1.0)]
    if region[3] - region[1] < 10.0:
        return None
    visual_rects = _collect_visual_rects(page, region)
    if not visual_rects:
        return None

    # caption x-범위와 어느 정도 겹치는 시각 객체를 우선 (열 정렬)
    cap_w = max(cap_bbox[2] - cap_bbox[0], 1.0)
    relevant = []
    for r in visual_rects:
        ov = _x_overlap(r, cap_bbox)
        if ov / cap_w > 0.05 or _x_overlap(r, [column_left, 0, column_right, 0]) > (column_right - column_left) * 0.05:
            relevant.append(r)
    if not relevant:
        relevant = visual_rects

    # caption에 가까운 cluster — y_bottom이 caption_top 근처(20pt 이내)에 있는 객체 우선
    cap_top = cap_bbox[1]
    relevant.sort(key=lambda r: cap_top - r[3])  # 가까운 거리(작은 양수)부터
    # 첫 객체를 시드로, y 방향으로 인접한 객체들을 cluster화
    cluster = [relevant[0]]
    cluster_top = relevant[0][1]
    cluster_bot = relevant[0][3]
    for r in relevant[1:]:
        # cluster의 vertical extent와 가깝거나 겹치면 포함
        gap = max(cluster_top - r[3], r[1] - cluster_bot)
        if gap <= 20.0:
            cluster.append(r)
            cluster_top = min(cluster_top, r[1])
            cluster_bot = max(cluster_bot, r[3])
    if not cluster:
        return None

    x0 = max(min(r[0] for r in cluster), region[0])
    y0 = max(min(r[1] for r in cluster), region[1])
    x1 = min(max(r[2] for r in cluster), region[2])
    y1 = min(max(r[3] for r in cluster), region[3])
    if x1 - x0 < 30.0 or y1 - y0 < 30.0:
        return None

    # 작은 패딩(figure 외곽선이 잘리지 않도록)
    pad = 4.0
    return [
        max(region[0], x0 - pad),
        max(region[1], y0 - pad),
        min(region[2], x1 + pad),
        min(region[3], y1 + pad),
    ]


def _check_contains_real_figure(page, bbox: List[float]) -> Tuple[bool, float]:
    """bbox 내부에 실제 figure visual content가 있는지 검사.

    Returns: (contains_real_figure, visual_density)
      - visual_density: visual rect 면적 / bbox 면적 (대략적인 시각 밀도)
    조건:
      - 시각 객체 면적 비율이 충분(>= 0.10)
      - 본문 텍스트 블록(>= 8자) 개수가 5개 미만
    """
    bbox_area = max((bbox[2] - bbox[0]) * (bbox[3] - bbox[1]), 1.0)
    vrects = _collect_visual_rects(page, bbox)
    if not vrects:
        return False, 0.0
    visual_area = 0.0
    for r in vrects:
        ix = max(0.0, min(bbox[2], r[2]) - max(bbox[0], r[0]))
        iy = max(0.0, min(bbox[3], r[3]) - max(bbox[1], r[1]))
        visual_area += ix * iy
    density = visual_area / bbox_area
    text_rects = _collect_text_rects(page, bbox, min_chars=20)
    if density < 0.10:
        return False, density
    if len(text_rects) >= 5 and density < 0.25:
        return False, density
    return True, density


def _detect_columns(page) -> List[Tuple[float, float]]:
    """페이지의 본문 컬럼(좌우 단) 추정. 텍스트 블록 x-중심을 보고 1열 또는 2열로 분류.
    Returns: [(left, right), ...] — 각 컬럼의 x 범위
    """
    page_w = page.rect.width
    try:
        blocks = page.get_text("dict").get("blocks", [])
    except Exception:
        return [(0.0, page_w)]
    centers = []
    for blk in blocks:
        if blk.get("type") != 0:
            continue
        bb = blk.get("bbox")
        if not bb:
            continue
        if bb[2] - bb[0] > page_w * 0.7:
            continue
        centers.append((bb[0] + bb[2]) / 2.0)
    if len(centers) < 4:
        return [(0.0, page_w)]
    # 페이지 중간을 기준으로 좌/우 분류
    mid = page_w / 2.0
    left_count = sum(1 for c in centers if c < mid)
    right_count = sum(1 for c in centers if c >= mid)
    if left_count >= 3 and right_count >= 3:
        return [(0.0, mid), (mid, page_w)]
    return [(0.0, page_w)]


def _column_for(bbox: List[float], columns: List[Tuple[float, float]]) -> Tuple[float, float]:
    """bbox 중심이 속한 컬럼 반환 (없으면 페이지 폭 전체)."""
    if not columns:
        return (0.0, 0.0)
    cx = (bbox[0] + bbox[2]) / 2.0
    for col in columns:
        if col[0] <= cx <= col[1]:
            return col
    return columns[0]


def _extract_figures_via_page_crop(doc: 'fitz.Document', existing_candidates: List[Dict],
                                    render_zoom: float = 2.0) -> List[Dict]:
    """이미지 추출 fallback (요청서 §6,§8,§10,§11).

    각 페이지에서 caption은 있는데 embedded image로 매칭되지 못한 경우,
    페이지를 렌더링해 caption_bbox 바로 위쪽 figure visual group 영역을 정밀하게 crop.

    개선:
      - 컬럼 인식 (좌/우 단 분리)
      - drawings + image bbox로 visual extent 정밀 계산
      - 본문 텍스트 영역은 crop에 포함하지 않음
      - contains_real_figure 검증 — text-only 영역 거부
    """
    import hashlib
    out: List[Dict] = []

    matched_fignum_by_page: Dict[int, set] = {}
    for c in existing_candidates:
        pg = c.get("page") or c.get("page_number")
        fn = _normalize_fig_number(c.get("figure_number","") or c.get("fig_label","")) or ""
        if pg and fn:
            matched_fignum_by_page.setdefault(pg, set()).add(fn)

    for pi in range(doc.page_count):
        try:
            page = doc.load_page(pi)
        except Exception:
            continue
        try:
            page_captions = _find_figure_captions_for_page(page)
        except Exception:
            page_captions = []
        if not page_captions:
            continue

        existing = matched_fignum_by_page.get(pi + 1, set())
        page_w = page.rect.width
        page_h = page.rect.height
        columns = _detect_columns(page)

        # 캡션을 y-top 순으로 정렬해 위쪽 cap 의 bottom을 다음 cap의 crop top 한계로 사용
        sorted_caps = sorted(page_captions, key=lambda c: c.get("bbox", [0,0,0,0])[1])

        for cap_idx, cap in enumerate(sorted_caps):
            fignum = _normalize_fig_number(cap.get("figure_number","")) or ""
            if not fignum:
                continue
            if fignum in existing:
                continue

            cap_bbox = cap.get("bbox")
            if not cap_bbox:
                continue

            # 컬럼 인식 — 같은 컬럼의 캡션만 upper_limit에 사용
            col_left, col_right = _column_for(cap_bbox, columns)
            if col_right - col_left < 50.0:
                col_left, col_right = 0.0, page_w

            upper_limit = 0.0
            for prev in sorted_caps[:cap_idx]:
                pb = prev.get("bbox")
                if not pb or pb[3] >= cap_bbox[1]:
                    continue
                # 같은 컬럼 캡션만 upper_limit으로 인정
                pcx = (pb[0] + pb[2]) / 2.0
                if not (col_left <= pcx <= col_right):
                    continue
                if pb[3] + 5.0 > upper_limit:
                    upper_limit = pb[3] + 5.0

            # 1) drawings + image bbox 기반 정밀 visual bbox
            visual_bbox = _compute_figure_visual_bbox(
                page, cap_bbox, upper_limit, col_left, col_right
            )

            crop_source = "visual_blocks_above_caption"
            contains_real_figure = True
            density = 0.0

            if visual_bbox is None:
                # 시각 객체 검출 실패 → 보수적 직사각형 fallback
                min_top_by_height = cap_bbox[1] - page_h * 0.55
                top = max(upper_limit, min_top_by_height, 0.0)
                bottom = max(top + 10.0, cap_bbox[1] - 2.0)
                pad_x = page_w * 0.04
                left = max(col_left, cap_bbox[0] - pad_x)
                right = min(col_right, cap_bbox[2] + pad_x)
                if (right - left) < (col_right - col_left) * 0.5:
                    left = col_left
                    right = col_right
                visual_bbox = [left, top, right, bottom]
                crop_source = "rect_fallback_above_caption"

            left, top, right, bottom = visual_bbox
            if (bottom - top) < 30.0 or (right - left) < 50.0:
                continue

            # text-only 영역 거부 검증 (요청서 §10)
            contains_real_figure, density = _check_contains_real_figure(page, visual_bbox)
            if not contains_real_figure:
                print(f"  [WARN] text-only candidate rejected: paper page={pi+1}, "
                      f"fig={cap.get('figure_number','?')}, density={density:.3f}")
                continue

            crop_rect = fitz.Rect(left, top, right, bottom)
            try:
                mat = fitz.Matrix(render_zoom, render_zoom)
                pix = page.get_pixmap(matrix=mat, clip=crop_rect, alpha=False)
                png_bytes = pix.tobytes("png")
            except Exception as e:
                print(f"  [WARN] page-crop render failed: page={pi+1}, fig={fignum}: {e}")
                continue

            try:
                png_bytes = trim_white_margin(png_bytes)
            except Exception:
                pass

            if Image:
                try:
                    im = Image.open(io.BytesIO(png_bytes))
                    img_w, img_h = im.size[0], im.size[1]
                except Exception:
                    img_w, img_h = pix.width, pix.height
            else:
                img_w, img_h = pix.width, pix.height

            # 너무 작은 crop은 스킵 (figure가 아닐 가능성)
            if img_w * img_h < 60_000:
                continue

            img_hash = hashlib.md5(png_bytes).hexdigest()
            _crop_cat, _crop_role, _crop_score = _classify_figure_category(cap["text"])
            _crop_fig_num = cap.get("figure_number", "")
            out.append({
                "page": pi + 1,
                "page_number": pi + 1,
                "figure_index": None,  # 외부에서 부여
                "area": img_w * img_h,
                "width": img_w,
                "height": img_h,
                "bbox": [left, top, right, bottom],
                "group_bbox": [left, top, right, bottom],
                "png_bytes": png_bytes,
                "byte_len": len(png_bytes),
                "caption_hint": cap["text"],
                "nearest_caption": cap["text"],
                "original_caption": cap["text"],
                "translated_caption": "",
                "caption_bbox": cap.get("bbox"),
                "caption_source": "page_crop_above_caption",
                "crop_source": crop_source,
                "contains_real_figure": True,
                "contains_body_text": False,
                "is_text_only": False,
                "visual_density": round(density, 3),
                "surrounding_text": "",
                "fig_label": fignum,
                "figure_number": _crop_fig_num,
                "image_figure_number": _crop_fig_num,
                "caption_figure_number": _crop_fig_num,  # page_crop은 caption 기반이므로 일치
                "category": _crop_cat,
                "role_candidate": _crop_role,
                "selection_score": round(_crop_score, 3),
                "hash": img_hash,
                "_img_hash": img_hash,
                "image_hash": img_hash,
                "image_source": "page_crop",
                "selection_reason": "embedded extraction missed; recovered via page-crop",
            })
            print(f"  [INFO] page-crop fallback figure created: page={pi+1}, "
                  f"fig={cap.get('figure_number','?')}, "
                  f"crop_source={crop_source}, density={density:.3f}, "
                  f"orig_caption='{cap['text'][:60]}'")

    return out


import base64


def _strip_invention_phrases(text: str) -> str:
    """캡션에서 '본 발명은/의/에', '본 연구의/에서', '본 논문의' 등 금지 표현 제거."""
    t = (text or "").strip()
    # "본 발명은", "본 발명의", "본 발명에 따른", "본 연구의", "본 연구에서", "본 논문의" 등 제거
    t = re.sub(r'본\s*(발명|연구|논문)(은|의|에서?|에\s*따른|에\s*의한|에서의)\s*', '', t)
    # 문두 공백/조사 정리
    t = re.sub(r'^\s*[,.]\s*', '', t).strip()
    return t


def _parse_fig_number_from_caption(caption: str) -> Optional[str]:
    """캡션 텍스트에서 Figure 번호를 추출.
    예: 'Fig. 3. Overview of ...' → '3'
        'Figure 12: Results' → '12'
        '그림 5 시스템 구조' → '5'
        'Fig.1(a) ...' → '1'
    """
    if not caption:
        return None
    # 다양한 Figure 표기 패턴
    m = re.search(r'(?:fig(?:ure)?|그림|도면)[.\s]*\s*(\d+)', caption, re.IGNORECASE)
    if m:
        return m.group(1)
    return None


def _normalize_fig_number(fig_num: str) -> Optional[str]:
    """fig_number 필드에서 순수 숫자만 추출.
    예: 'Figure 3' → '3', 'Fig.12' → '12', '5' → '5', 'fig1' → '1'
    """
    if not fig_num:
        return None
    m = re.search(r'(\d+)', str(fig_num))
    return m.group(1) if m else None


def normalize_caption(caption: str) -> str:
    """캡션 표시 전 정규화 — Fig.Fig. / Figure Figure / 그림 그림 같은 중복 prefix 제거.
    요청서 §10 규칙."""
    if not caption:
        return ""
    c = str(caption).strip()
    # 'Fig.Fig.', 'Fig. Fig.', 'Fig.  Fig.' 등 중복 → 'Fig. '
    c = re.sub(r"^(?:Fig\.?\s*){2,}", "Fig. ", c, flags=re.IGNORECASE)
    c = re.sub(r"^(?:Figure\.?\s*){2,}", "Figure ", c, flags=re.IGNORECASE)
    c = re.sub(r"^(?:그림\s*){2,}", "그림 ", c)
    c = re.sub(r"^(?:도면\s*){2,}", "도면 ", c)
    # 'Fig.' 바로 뒤가 'Figure'인 경우
    c = re.sub(r"^Fig\.?\s*Figure\s*", "Fig. ", c, flags=re.IGNORECASE)
    c = re.sub(r"\s+", " ", c).strip()
    return c


def _match_candidate_by_fig_info(
    candidates: List[Dict],
    fig_number: str,
    fig_page: Optional[int] = None,
) -> Optional[int]:
    """텍스트 LLM이 선정한 fig_number/page에 매칭되는 후보 이미지의 인덱스를 반환.
    
    매칭 전략:
    0) fig_label 필드 직접 비교 (캡션 파싱으로 추출해둔 번호) — 가장 정확
    1) nearest_caption에서 Figure 번호를 파싱하여 fig_number와 비교
    2) page 번호도 일치하면 우선
    3) page가 다르더라도 fig_number만 일치하면 후보로
    """
    target_num = _normalize_fig_number(fig_number)
    if not target_num:
        return None

    # 0차: fig_label 직접 비교 + page 일치 (가장 신뢰도 높음)
    if fig_page is not None:
        for idx, c in enumerate(candidates):
            if _normalize_fig_number(c.get("fig_label", "")) == target_num and c.get("page") == fig_page:
                print(f"    [MATCH-0] fig_label='{c.get('fig_label')}' page={fig_page} → candidate[{idx}]")
                return idx

    # 0차-B: fig_label 직접 비교 (page 무관)
    for idx, c in enumerate(candidates):
        if _normalize_fig_number(c.get("fig_label", "")) == target_num:
            print(f"    [MATCH-0B] fig_label='{c.get('fig_label')}' → candidate[{idx}] (page={c.get('page')})")
            return idx
    
    # 1차: page + nearest_caption fig_number 모두 일치
    if fig_page is not None:
        for idx, c in enumerate(candidates):
            cap = c.get("nearest_caption", "") or c.get("caption_hint", "")
            parsed_num = _parse_fig_number_from_caption(cap)
            if parsed_num == target_num and c.get("page") == fig_page:
                return idx
    
    # 2차: fig_number만 일치 (page 불일치 허용)
    for idx, c in enumerate(candidates):
        cap = c.get("nearest_caption", "") or c.get("caption_hint", "")
        parsed_num = _parse_fig_number_from_caption(cap)
        if parsed_num == target_num:
            return idx
    
    # 3차: caption_hint (페이지 전체 캡션)에서 fig_number 검색
    for idx, c in enumerate(candidates):
        hint = c.get("caption_hint", "")
        if not hint:
            continue
        for part in hint.split(" | "):
            parsed_num = _parse_fig_number_from_caption(part)
            if parsed_num == target_num:
                if fig_page is None or c.get("page") == fig_page:
                    return idx
    
    # 4차: page만 일치하는 후보 중 가장 큰 이미지 (area 기준, 이미 정렬됨)
    if fig_page is not None:
        for idx, c in enumerate(candidates):
            if c.get("page") == fig_page:
                return idx
    
    return None


def _translate_captions_only(
    cap1_raw: str,
    cap2_raw: str,
    fig_title_context: str = "",
    model: str = DEFAULT_MODEL,
) -> List[str]:
    """이미지가 이미 확정된 경우, 원문 캡션을 한국어로 번역만 수행.
    Returns: [concept_caption_ko, result_caption_ko]
    실패 시 원문 캡션 그대로 반환.
    """
    api_key = os.getenv("OPENAI_API_KEY", "").strip()
    gemini_key = os.getenv("GEMINI_API_KEY", "").strip()

    prompt = (
        "아래 두 Figure 원문 캡션(PDF에서 직접 추출됨)을 각각 한국어로 번역하세요.\n"
        "번역 규칙 (반드시 준수):\n"
        "- **반드시 제공된 원문 캡션 텍스트만 번역할 것. 새로 만들거나 추측 금지.**\n"
        "- **Fig. 번호 (Fig. 1., Figure 2., 그림 3 등)는 번역 결과에서 그대로 유지할 것.**\n"
        "  예: 'Fig. 1. Circuit diagram of 3ph LCL-filtered VSC IOBC system'\n"
        "      → 'Fig. 1. 3상 LCL 필터 VSC IOBC 시스템 회로도'\n"
        "- 명사로 끝낼 것\n"
        "- '본 발명은/의', '본 연구의', '본 논문의' 등 표현 금지\n"
        "- 효과·결과 서술 금지 (무엇을 보여주는 그림인지만 설명)\n"
        "- 숫자+단위 붙여쓰기 (50kHz, 1.2A 등)\n"
        "- '그림1', '그림2', 'Figure 1', 'Figure 2'처럼 번호만 적힌 라벨은 캡션으로 사용하지 말 것 — 반드시 그림 내용 설명 포함\n"
        "- **[중요] 원문 캡션이 '(캡션 없음)' 이거나 비어있으면, 임의로 캡션을 만들지 말고 해당 caption은 빈 문자열(\"\")로 반환할 것**\n"
        + (f"\n{fig_title_context}" if fig_title_context else "") +
        f"\n\nFigure 1 원문 캡션: {cap1_raw or '(캡션 없음)'}\n"
        f"Figure 2 원문 캡션: {cap2_raw or '(캡션 없음)'}\n\n"
        "출력 형식 (JSON):\n"
        '{"caption1": "번역된 캡션1 (Fig. 번호 유지, 원문 없으면 빈 문자열)", "caption2": "번역된 캡션2 (Fig. 번호 유지, 원문 없으면 빈 문자열)"}'
    )
    try:
        if IS_GEMINI_MODEL(model) and HAS_GEMINI and gemini_key:
            import google.generativeai as genai
            genai.configure(api_key=gemini_key)
            # response_mime_type\uc740 gemini-2.5-preview \uad6c\ubc84\uc804 \ud638\ud658\uc131 \ubb38\uc81c\ub85c \uc81c\uac70,
            # \ub300\uc2e0 JSON \ucf54\ub4dc\ube14\ub85d \ud30c\uc2f1 fallback \uc801\uc6a9
            gm = genai.GenerativeModel(model)
            resp = gm.generate_content(prompt + "\n\n\ucd9c\ub825\uc740 \ubc18\ub4dc\uc2dc JSON\ub9cc \ucd9c\ub825. \ucf54\ub4dc\ube14\ub85d \uc5c6\uc774.")
            raw = (resp.text or "").strip()
            # JSON 코드블록 없애주기
            jm = re.search(r"```(?:json)?\s*([\s\S]+?)\s*```", raw)
            parsed = json.loads(jm.group(1).strip() if jm else raw)
        elif api_key:
            client = OpenAI(api_key=api_key)
            is_reasoning = any(x in model.lower() for x in ["o1", "o3", "gpt-5", "gpt-o3"])
            params: Dict[str, Any] = {
                "model": model,
                "messages": [{"role": "user", "content": prompt}],
            }
            if is_reasoning:
                params["max_completion_tokens"] = 500
            else:
                params["response_format"] = {"type": "json_object"}
                params["max_completion_tokens"] = 300
                params["temperature"] = 0.3
            resp = client.chat.completions.create(**params)
            content = (resp.choices[0].message.content or "").strip()
            jm = re.search(r"```(?:json)?\s*([\s\S]+?)\s*```", content)
            parsed = json.loads(jm.group(1) if jm else content)
        else:
            raise RuntimeError("No API key")
        c1 = parsed.get("caption1", "") or ""
        c2 = parsed.get("caption2", "") or ""
        print(f"  [INFO] Caption translated: '{c1[:30]}' / '{c2[:30]}'")
        return [c1, c2]
    except Exception as e:
        print(f"  [WARN] Caption translation failed: {e}. Returning empty captions (will use fig_title fallback).")
        # \uc601\uc5b4 \uc6d0\ubb38\uc744 \uadf8\ub300\ub85c \ubc18\ud658\ud558\uc9c0 \uc54a\uc74c \u2014 normalize_brief \ub2e8\uacc4\uc5d0\uc11c fig_title \uae30\ubc18\uc73c\ub85c \uc7ac\ucc98\ub9ac\ub428
        return ["", ""]


def _figures_are_duplicate(a: Dict, b: Dict) -> Tuple[bool, str]:
    """대표 이미지 2개가 사실상 동일한 이미지인지 다중 기준으로 판별.
    Returns: (is_duplicate, reason)
    기준 (요청서 §7):
      1) 이미지 바이트 해시 일치
      2) 동일 figure_index
      3) 같은 페이지 + bbox overlap > 0.7
      4) 동일 width × height + 바이트 길이 차이 5% 이하
      5) 동일 figure_number(또는 fig_label)
      6) 동일 original_caption (정규화 후)
    """
    if not a or not b:
        return False, ""
    import hashlib
    a_hash = a.get("hash") or a.get("_img_hash") or a.get("image_hash")
    b_hash = b.get("hash") or b.get("_img_hash") or b.get("image_hash")
    if not a_hash and a.get("png_bytes"):
        a_hash = hashlib.md5(a["png_bytes"]).hexdigest()
    if not b_hash and b.get("png_bytes"):
        b_hash = hashlib.md5(b["png_bytes"]).hexdigest()
    if a_hash and b_hash and a_hash == b_hash:
        return True, f"identical_hash({a_hash[:8]})"

    if (a.get("figure_index") is not None
            and a.get("figure_index") == b.get("figure_index")):
        return True, f"identical_figure_index({a.get('figure_index')})"

    a_page = a.get("page") or a.get("page_number")
    b_page = b.get("page") or b.get("page_number")
    a_bbox = a.get("bbox")
    b_bbox = b.get("bbox")
    if a_page is not None and a_page == b_page and a_bbox and b_bbox:
        try:
            ratio = _compute_overlap_ratio(a_bbox, b_bbox)
        except Exception:
            ratio = 0.0
        if ratio > 0.7:
            return True, f"same_page_bbox_overlap({ratio:.2f})"

    a_w, a_h = a.get("width"), a.get("height")
    b_w, b_h = b.get("width"), b.get("height")
    a_bl, b_bl = a.get("byte_len"), b.get("byte_len")
    if a_bl is None and a.get("png_bytes"):
        a_bl = len(a["png_bytes"])
    if b_bl is None and b.get("png_bytes"):
        b_bl = len(b["png_bytes"])
    if a_w and b_w and a_h and b_h and a_bl and b_bl:
        if a_w == b_w and a_h == b_h:
            denom = max(a_bl, b_bl)
            if denom and abs(a_bl - b_bl) / float(denom) < 0.05:
                return True, f"same_dims_close_bytes({a_w}x{a_h})"

    # 5) 동일 figure_number / fig_label
    a_fn = _normalize_fig_number(a.get("figure_number","") or a.get("fig_label","")) or ""
    b_fn = _normalize_fig_number(b.get("figure_number","") or b.get("fig_label","")) or ""
    if a_fn and b_fn and a_fn == b_fn:
        # 같은 페이지에서 같은 figure_number는 거의 확실히 같은 figure
        if a_page is not None and a_page == b_page:
            return True, f"same_figure_number_same_page(Fig.{a_fn} p{a_page})"
        # 다른 페이지여도 동일 fig_number는 중복으로 본다(논문 1편 기준)
        return True, f"same_figure_number(Fig.{a_fn})"

    # 6) 동일 original_caption
    def _norm_cap(t: str) -> str:
        if not t:
            return ""
        x = re.sub(r"\s+", " ", str(t)).strip().lower()
        return x
    a_cap = _norm_cap(a.get("original_caption","") or a.get("nearest_caption",""))
    b_cap = _norm_cap(b.get("original_caption","") or b.get("nearest_caption",""))
    if a_cap and b_cap and a_cap == b_cap and len(a_cap) > 15:
        return True, f"identical_original_caption('{a_cap[:30]}...')"

    return False, ""


def _find_non_duplicate_alt(top_candidates: List[Dict], exclude_idx: int,
                            reference_fig: Dict) -> Optional[int]:
    """top_candidates에서 reference_fig와 중복이 아닌 가장 큰 후보의 인덱스를 찾는다."""
    for alt_idx in range(len(top_candidates)):
        if alt_idx == exclude_idx:
            continue
        is_dup, _ = _figures_are_duplicate(top_candidates[alt_idx], reference_fig)
        if not is_dup:
            return alt_idx
    return None


def _figure_log_summary(label: str, fig: Optional[Dict], cap_ko: str = "",
                        selection_reason: str = "") -> str:
    """선택된 figure 1개에 대한 한 줄 로그 문자열 — page/figure_index/fig_label/hash/caption/reason."""
    if not fig:
        return f"  [SELECTED] {label}: <None>  reason={selection_reason or '?'}"
    h = fig.get("hash") or fig.get("_img_hash") or ""
    return (
        f"  [SELECTED] {label}: "
        f"page={fig.get('page') or fig.get('page_number') or '?'} "
        f"figure_index={fig.get('figure_index', '?')} "
        f"fig_label='{fig.get('fig_label', '')}' "
        f"hash={h[:8]} "
        f"size={fig.get('width','?')}x{fig.get('height','?')} "
        f"caption_source='{fig.get('caption_source','')}' "
        f"orig_caption='{(fig.get('nearest_caption') or '')[:60]}' "
        f"final_caption='{(cap_ko or '')[:60]}' "
        f"reason='{selection_reason}'"
    )


def pick_two_figures_with_vision(candidates: List[Dict], paper_text: str, model: str = DEFAULT_MODEL, representative_figures: Optional[List[Dict]] = None) -> Tuple[Optional[Dict], Optional[Dict], List[str]]:
    """
    Use Vision LLM to pick the best concept and result images, and generate their captions.
    representative_figures가 제공되면 fig_title을 캡션 생성의 기준으로 활용.
    Returns: (concept_fig, effect_fig, [caption1, caption2])
    - 이미지를 아예 추출 못한 경우 → (None, None, ["", ""])
    - 이미지를 선택했으나 png_bytes 없으면 fig_number 캡션만 반환
    """
    if not candidates:
        return None, None, ["", ""]

    # ── fig_title_context: 캡션 생성 컨텍스트 (매칭 스킵 경로에서도 사용) ──
    fig_title_context = ""
    if representative_figures and len(representative_figures) >= 2:
        fig1 = representative_figures[0]
        fig2 = representative_figures[1]
        fig_title_context = (
            "\n\n**[중요] 텍스트 분석 결과 선정된 대표 이미지 정보:**\n"
            f"- 컨셉 이미지: Fig.{fig1.get('fig_number', '?')} — \"{fig1.get('fig_title', '')}\""
            f" (page {fig1.get('page', '?')})\n"
            f"- 결과 이미지: Fig.{fig2.get('fig_number', '?')} — \"{fig2.get('fig_title', '')}\""
            f" (page {fig2.get('page', '?')})\n"
            "위 정보를 참고하여, 해당 Figure와 가장 일치하는 이미지를 선택하고, "
            "원문 Figure 제목(fig_title)을 기반으로 한국어 캡션을 작성하세요.\n"
        )

    # ── 텍스트 LLM의 representative_figures로 후보 매칭 시도 ──
    matched_concept_idx = None
    matched_result_idx = None
    if representative_figures and len(representative_figures) >= 2:
        fig1 = representative_figures[0]  # concept
        fig2 = representative_figures[1]  # result
        matched_concept_idx = _match_candidate_by_fig_info(
            candidates, fig1.get("fig_number", ""), fig1.get("page")
        )
        matched_result_idx = _match_candidate_by_fig_info(
            candidates, fig2.get("fig_number", ""), fig2.get("page")
        )
        # 같은 후보에 매칭되면 result 매칭 무효화
        if matched_concept_idx is not None and matched_result_idx is not None:
            if matched_concept_idx == matched_result_idx:
                print(f"  [WARN] Both figures matched to same candidate idx={matched_concept_idx}. Clearing result match.")
                matched_result_idx = None
        
        if matched_concept_idx is not None and matched_result_idx is not None:
            print(f"  [INFO] Figure matching: concept → candidate[{matched_concept_idx}] "
                  f"(Fig.{fig1.get('fig_number','?')}, p{fig1.get('page','?')}), "
                  f"result → candidate[{matched_result_idx}] "
                  f"(Fig.{fig2.get('fig_number','?')}, p{fig2.get('page','?')})")
        else:
            c_status = f"candidate[{matched_concept_idx}]" if matched_concept_idx is not None else "NOT FOUND"
            r_status = f"candidate[{matched_result_idx}]" if matched_result_idx is not None else "NOT FOUND"
            print(f"  [WARN] Partial figure match: concept → {c_status}, result → {r_status}. "
                  f"Vision LLM will select unmatched figures.")

    # Take top 6 largest images to avoid huge payloads
    top_candidates = candidates[:6]

    # 매칭된 후보가 top_candidates 범위 밖이면 강제 포함
    for m_idx in [matched_concept_idx, matched_result_idx]:
        if m_idx is not None and m_idx >= len(top_candidates):
            top_candidates.append(candidates[m_idx])
            print(f"  [INFO] Matched candidate[{m_idx}] was outside top-6, added to candidates list.")

    # 매칭된 인덱스를 top_candidates 기준으로 재계산
    if matched_concept_idx is not None and matched_concept_idx >= 6:
        for i, c in enumerate(top_candidates):
            if c is candidates[matched_concept_idx]:
                matched_concept_idx = i
                break
    if matched_result_idx is not None and matched_result_idx >= 6:
        for i, c in enumerate(top_candidates):
            if c is candidates[matched_result_idx]:
                matched_result_idx = i
                break

    # ── 두 이미지 모두 fig_label로 완전 매칭 → Vision LLM 스킵 ──
    # → 이미지 선택 오류 가능성을 원천 차단
    if matched_concept_idx is not None and matched_result_idx is not None:
        cf = candidates[matched_concept_idx] if matched_concept_idx < len(candidates) else top_candidates[matched_concept_idx]
        ef = candidates[matched_result_idx] if matched_result_idx < len(candidates) else top_candidates[matched_result_idx]

        # 매칭 결과가 사실상 동일 이미지인 경우 → ef를 중복 아닌 후보로 교체
        _is_dup, _dup_reason = _figures_are_duplicate(cf, ef)
        if _is_dup:
            print(f"  [WARN] Matched concept/result are duplicates ({_dup_reason}). Searching alternative...")
            for _ai in range(len(candidates)):
                if _ai == matched_concept_idx:
                    continue
                _alt_dup, _ = _figures_are_duplicate(cf, candidates[_ai])
                if not _alt_dup:
                    ef = candidates[_ai]
                    matched_result_idx = _ai
                    print(f"    -> Replaced result with candidate[{_ai}] "
                          f"(page={ef.get('page','?')}, hash={(ef.get('hash') or '')[:8]})")
                    break

        print(f"  [INFO] Both figures fully matched\u2192 skipping Vision LLM image selection.")

        # 원문 캡션을 한국어로 번역 — 반드시 PDF 원문(original_caption)만 사용 (요청서 §3-2)
        cap1_raw = cf.get("original_caption", "") or cf.get("nearest_caption", "")
        cap2_raw = ef.get("original_caption", "") or ef.get("nearest_caption", "")
        translated = _translate_captions_only(cap1_raw, cap2_raw,
                                              fig_title_context=fig_title_context,
                                              model=model)
        concept_cap = _strip_invention_phrases(translated[0])
        result_cap  = _strip_invention_phrases(translated[1])
        # 선택된 figure dict에 번역 결과 저장 (PPT 표시 단계에서 우선 사용)
        cf["translated_caption"] = concept_cap
        ef["translated_caption"] = result_cap
        cf["selection_reason"] = "text-LLM fig_number/page match"
        ef["selection_reason"] = "text-LLM fig_number/page match"
        # ── 최종 선택 결과 로그 (matched-skip 경로) ──
        print(_figure_log_summary("concept", cf, concept_cap,
                                  "text-LLM fig_number/page match (Vision LLM skipped)"))
        print(_figure_log_summary("result",  ef, result_cap,
                                  "text-LLM fig_number/page match (Vision LLM skipped)"))
        if not cf.get("png_bytes"):
            cf = None
        if not ef.get("png_bytes"):
            ef = None
        return cf, ef, [concept_cap, result_cap]

    api_key = os.getenv("OPENAI_API_KEY", "").strip()
    gemini_key = os.getenv("GEMINI_API_KEY", "").strip()

    # (fig_title_context는 앞서 이미 정의됨)

    CAPTION_INSTRUCTION = (
        "PDF에서 추출된 figure 후보 리스트가 제공됩니다.\n"
        "각 후보에는 image_path, page_number, figure_number, original_caption, category, metadata가 포함됩니다.\n\n"
        "대표 figure는 정확히 2개만 선택하세요.\n\n"
        "1. concept_index (top_image):\n"
        "제안 기술의 구체적인 방법, 기술 컨셉, 구조, 아키텍처, 회로, 모델, 제어 구조, 시스템 구성을 가장 잘 설명하는 figure를 선택하세요.\n"
        "category='method_or_concept'인 후보를 우선 선택하세요.\n\n"
        "2. result_index (bottom_image):\n"
        "개선효과, 실험 결과, 검증 결과, 성능 비교, 효율, 오차, 파형, 평가 결과를 가장 잘 보여주는 figure를 선택하세요.\n"
        "category='result_or_effect'인 후보를 우선 선택하세요.\n\n"
        "**[매우 중요] 선택 규칙:**\n"
        "- concept_index와 result_index는 반드시 서로 다른 번호여야 합니다.\n"
        "- 중복 figure를 선택하지 마세요.\n"
        "- 텍스트만 있는 후보는 선택하지 마세요.\n"
        "- figure_number, original_caption, group_bbox, image_hash가 같은 후보를 중복 선택하지 마세요.\n"
        "- caption을 추측하지 마세요. 반드시 제공된 original_caption만 사용하세요.\n"
        "- 선택된 이미지와 caption은 반드시 같은 figure_number여야 합니다.\n"
        "- image_figure_number와 caption_figure_number가 다르면 해당 후보는 제외하세요.\n"
        "- original_caption이 비어있는 후보보다 캡션이 정상적으로 추출된 후보를 우선 선택하세요.\n\n"
        + fig_title_context +
        "\n출력은 반드시 다음 JSON 형식으로 하세요 (마크다운 코드블록 없이):\n"
        '{\n  "concept_index": 0,\n  "result_index": 1\n}'
    )

    content = [{"type": "text", "text": CAPTION_INSTRUCTION}]
    context_text = paper_text[:5000]
    if len(paper_text) > 5000:
        context_text += "\n...\n" + paper_text[5000:8000]
    content.append({"type": "text", "text": f"\n\n[문서 내용 일부]\n{context_text}"})

    # Add images with caption hints, figure_number, image_hash, category and surrounding text
    for idx, c in enumerate(top_candidates):
        b64 = base64.b64encode(c["png_bytes"]).decode("utf-8")
        cap_info = c.get("original_caption") or c.get("nearest_caption") or c.get("caption_hint", "")
        fignum = c.get("figure_number", "") or (f"Fig. {c.get('fig_label')}" if c.get('fig_label') else "")
        ihash = (c.get("image_hash") or c.get("hash") or c.get("_img_hash") or "")[:8]
        cat = c.get("category", "other")
        role = c.get("role_candidate", "")
        meta_parts = [f"figure_number={fignum or '?'}", f"image_hash={ihash or '?'}",
                      f"page={c.get('page','?')}", f"category={cat}", f"role={role or '?'}"]
        meta_str = " | ".join(meta_parts)
        cap_info_str = f"\n  [original_caption]: {cap_info}" if cap_info else "\n  [original_caption]: (없음)"
        surrounding = c.get("surrounding_text", "")
        surrounding_str = f"\n  [주변 문맥]: {surrounding[:500]}" if surrounding else ""
        content.append({
            "type": "text",
            "text": f"--- 이미지 후보 번호: {idx} | {meta_str}{cap_info_str}{surrounding_str} ---"
        })
        content.append({
            "type": "image_url",
            "image_url": {"url": f"data:image/png;base64,{b64}", "detail": "low"}
        })

    def _call_vision_openai(client, model):
        is_reasoning_model = any(x in model.lower() for x in ["o1", "o3", "gpt-5", "gpt-o3"])
        params: Dict[str, Any] = {
            "model": model,
            "messages": [{"role": "user", "content": content}],
        }
        if is_reasoning_model:
            params["max_completion_tokens"] = 4000
        else:
            params["response_format"] = {"type": "json_object"}
            params["max_completion_tokens"] = 600
        resp = client.chat.completions.create(**params)
        return (resp.choices[0].message.content or "").strip()

    try:
        if IS_GEMINI_MODEL(model) and HAS_GEMINI and gemini_key:
            # Gemini vision: PIL Image 객체로 변환 (SDK가 가장 안정적으로 처리)
            import google.generativeai as genai
            from PIL import Image as _PILImage
            import io as _io_cap
            genai.configure(api_key=gemini_key)
            gmodel = genai.GenerativeModel(model)
            # 컨텍스트를 OpenAI와 동일하게 8000자까지 확대
            context_text = paper_text[:5000]
            if len(paper_text) > 5000:
                context_text += "\n...\n" + paper_text[5000:8000]
            parts = [
                CAPTION_INSTRUCTION,
                f"\n[문서 내용 일부]\n{context_text}"
            ]
            for idx, c in enumerate(top_candidates):
                cap_info = c.get("original_caption") or c.get("nearest_caption") or c.get("caption_hint", "")
                fignum = c.get("figure_number", "") or (f"Fig. {c.get('fig_label')}" if c.get('fig_label') else "")
                ihash = (c.get("image_hash") or c.get("hash") or c.get("_img_hash") or "")[:8]
                cat = c.get("category", "other")
                role = c.get("role_candidate", "")
                meta_str = f"figure_number={fignum or '?'} | image_hash={ihash or '?'} | page={c.get('page','?')} | category={cat} | role={role or '?'}"
                cap_str = f"\n  [original_caption]: {cap_info}" if cap_info else "\n  [original_caption]: (없음)"
                surrounding = c.get("surrounding_text", "")
                surrounding_str = f"\n  [주변 문맥]: {surrounding[:500]}" if surrounding else ""
                parts.append(f"--- 이미지 후보 번호: {idx} | {meta_str}{cap_str}{surrounding_str} ---")
                # PIL Image로 변환하여 전달 (Gemini SDK 호환성 최적)
                pil_img = _PILImage.open(_io_cap.BytesIO(c["png_bytes"]))
                parts.append(pil_img)
            resp = gmodel.generate_content(parts)
            res_text = (resp.text or "").strip()
        elif api_key:
            client = OpenAI(api_key=api_key)
            res_text = _call_vision_openai(client, model)
        else:
            raise RuntimeError("No API key available for vision call")

        json_match = re.search(r"```(?:json)?\s*([\s\S]+?)\s*```", res_text)
        if json_match:
            res_text = json_match.group(1).strip()

        parsed = json.loads(res_text)
        c_idx = int(parsed.get("concept_index", 0))
        r_idx = int(parsed.get("result_index", 1))

        c_idx = c_idx if 0 <= c_idx < len(top_candidates) else 0
        r_idx = r_idx if 0 <= r_idx < len(top_candidates) else min(1, len(top_candidates) - 1)

        # ── 텍스트 LLM 매칭 결과를 Vision LLM 선택보다 우선 적용 ──
        if matched_concept_idx is not None:
            if c_idx != matched_concept_idx:
                print(f"  [INFO] Overriding Vision LLM concept selection: "
                      f"idx {c_idx} → {matched_concept_idx} (text LLM match)")
            c_idx = matched_concept_idx
        if matched_result_idx is not None:
            if r_idx != matched_result_idx:
                print(f"  [INFO] Overriding Vision LLM result selection: "
                      f"idx {r_idx} → {matched_result_idx} (text LLM match)")
            r_idx = matched_result_idx

        # 중복 방지 1단계: 인덱스가 같은 경우
        if c_idx == r_idx and len(top_candidates) > 1:
            print(f"  [WARN] Vision LLM selected same index for concept and result (idx={c_idx}). Auto-fixing...")
            for alt_idx in range(len(top_candidates)):
                if alt_idx != c_idx:
                    r_idx = alt_idx
                    break

        # 중복 방지 2단계: 다중 기준 중복 검사 (hash / bbox / 크기·바이트 / figure_index)
        is_dup, dup_reason = _figures_are_duplicate(top_candidates[c_idx], top_candidates[r_idx])
        if is_dup and len(top_candidates) > 1:
            print(f"  [WARN] Vision LLM selected duplicate figures "
                  f"(c_idx={c_idx}, r_idx={r_idx}, reason={dup_reason}). Auto-fixing...")
            alt_idx = _find_non_duplicate_alt(top_candidates, c_idx, top_candidates[c_idx])
            if alt_idx is not None:
                r_idx = alt_idx
                alt_hash = (top_candidates[alt_idx].get("hash") or
                            top_candidates[alt_idx].get("_img_hash") or "")
                print(f"    -> Replaced result with candidate[{alt_idx}] "
                      f"(page={top_candidates[alt_idx].get('page','?')}, hash={alt_hash[:8]})")
            else:
                print(f"    -> No non-duplicate alternative found in top_candidates.")

        concept_fig = top_candidates[c_idx]
        effect_fig = top_candidates[r_idx]
        print(f"  [INFO] Final image selection: concept=candidate[{c_idx}] "
              f"(page {concept_fig.get('page','?')}, cap='{(concept_fig.get('nearest_caption','') or '')[:40]}'), "
              f"result=candidate[{r_idx}] "
              f"(page {effect_fig.get('page','?')}, cap='{(effect_fig.get('nearest_caption','') or '')[:40]}')")

        # 이미지 바이트가 없으면 None 처리, 캡션만 보존
        def _fig_or_none(fig, fallback_cap):
            if not fig or not fig.get("png_bytes"):
                return None, fallback_cap
            return fig, fallback_cap

        # §10: Vision LLM은 선택만 수행. 캡션 번역은 최종 선택된 2개에 대해서만 수행
        cap1_raw = concept_fig.get("original_caption", "") or concept_fig.get("nearest_caption", "")
        cap2_raw = effect_fig.get("original_caption", "") or effect_fig.get("nearest_caption", "")
        if cap1_raw or cap2_raw:
            translated = _translate_captions_only(cap1_raw, cap2_raw,
                                                  fig_title_context=fig_title_context,
                                                  model=model)
            concept_cap = _strip_invention_phrases(translated[0])
            result_cap = _strip_invention_phrases(translated[1])
        else:
            concept_cap = ""
            result_cap = ""

        # ── 최종 선택 사유 결정 ──
        c_reason = "Vision LLM concept selection"
        if matched_concept_idx is not None:
            c_reason += " (overridden by text-LLM fig match)"
        r_reason = "Vision LLM result selection"
        if matched_result_idx is not None:
            r_reason += " (overridden by text-LLM fig match)"

        # ── 번역 결과를 figure dict에 보존 (PPT 표시 단계에서 우선 사용) ──
        if concept_fig is not None:
            concept_fig["translated_caption"] = concept_cap
            concept_fig["selection_reason"] = c_reason
        if effect_fig is not None:
            effect_fig["translated_caption"] = result_cap
            effect_fig["selection_reason"] = r_reason

        concept_fig, concept_cap = _fig_or_none(concept_fig, concept_cap)
        effect_fig, result_cap = _fig_or_none(effect_fig, result_cap)
        print(_figure_log_summary("concept", concept_fig, concept_cap, c_reason))
        print(_figure_log_summary("result",  effect_fig,  result_cap,  r_reason))

        return concept_fig, effect_fig, [concept_cap, result_cap]

    except Exception as e:
        print(f"  [WARN] Vision API failed: {e}. Falling back to heuristic.")
        try:
            c, e_fig = pick_two_figures(candidates)
            # heuristic 폴백 시에도 PDF 원문 캡션이 있으면 번역 시도
            cap1_raw = (c.get("original_caption", "") or c.get("nearest_caption", "")) if c else ""
            cap2_raw = (e_fig.get("original_caption", "") or e_fig.get("nearest_caption", "")) if e_fig else ""
            if cap1_raw or cap2_raw:
                try:
                    translated = _translate_captions_only(
                        cap1_raw, cap2_raw,
                        fig_title_context=fig_title_context,
                        model=model,
                    )
                    cap1 = _strip_invention_phrases(translated[0])
                    cap2 = _strip_invention_phrases(translated[1])
                except Exception:
                    cap1, cap2 = "", ""
            else:
                cap1, cap2 = "", ""
            if c is not None:
                c["translated_caption"] = cap1
                c["selection_reason"] = "heuristic fallback (Vision LLM failed)"
            if e_fig is not None:
                e_fig["translated_caption"] = cap2
                e_fig["selection_reason"] = "heuristic fallback (Vision LLM failed)"
            return c, e_fig, [cap1, cap2]
        except Exception:
            return None, None, ["", ""]

def pick_two_figures(candidates: List[Dict]) -> Tuple[Dict, Dict]:
    """Category 기반 역할별 figure 선택 (요청서 §5).
    top_image: method_or_concept category에서 가장 대표적인 figure
    bottom_image: result_or_effect category에서 가장 대표적인 figure
    """
    import hashlib
    if not candidates:
        raise RuntimeError("No image candidates found in PDF.")

    def _get_hash(c):
        return c.get("_img_hash") or c.get("image_hash") or hashlib.md5(c["png_bytes"]).hexdigest()

    # Category 기반 분리
    method_cands = [(i, c) for i, c in enumerate(candidates) if c.get("category") == "method_or_concept"]
    result_cands = [(i, c) for i, c in enumerate(candidates) if c.get("category") == "result_or_effect"]

    # method_or_concept에서 selection_score + area 기준 최고점
    ci = None
    if method_cands:
        method_cands.sort(key=lambda x: (x[1].get("selection_score", 0), x[1].get("area", 0)), reverse=True)
        ci = method_cands[0][0]

    # result_or_effect에서 selection_score + area 기준 최고점 (ci와 다른 것)
    ei = None
    if result_cands:
        result_cands.sort(key=lambda x: (x[1].get("selection_score", 0), x[1].get("area", 0)), reverse=True)
        for idx, c in result_cands:
            if idx != ci:
                ei = idx
                break

    # Fallback: category 기반 선택 실패 시 기존 키워드 방식
    if ci is None or ei is None:
        concept_kw = ["framework", "architecture", "pipeline", "overview", "system",
                      "method", "approach", "proposed", "module", "flow", "diagram"]
        effect_kw = ["result", "performance", "evaluation", "comparison", "error",
                     "accuracy", "rmse", "improvement", "reduction", "%"]

        def score(c, kws):
            t = (c.get("caption_hint", "") or c.get("original_caption", "")).lower()
            return sum(1 for k in kws if k in t)

        scored = [(i, score(c, concept_kw), score(c, effect_kw), c["area"]) for i, c in enumerate(candidates)]

        if ci is None:
            ci = sorted(scored, key=lambda x: (x[1], x[3]), reverse=True)[0][0]

        c_hash = _get_hash(candidates[ci])

        if ei is None:
            for i, _, _, _ in sorted(scored, key=lambda x: (x[2], x[3]), reverse=True):
                if i != ci and _get_hash(candidates[i]) != c_hash:
                    ei = i
                    break

    if ei is None:
        for i in range(len(candidates)):
            if i != ci:
                ei = i
                break
    if ei is None:
        ei = ci  # 후보가 1개뿐
    return candidates[ci], candidates[ei]


# ============================================================
# WIPS Patent URL Scraping
# ============================================================
def scrape_wips_patent(url: str) -> Tuple[str, str]:
    """
    WIPS 특허 페이지에서 초록, 청구항, 발명 설명 등 텍스트를 스크래핑.
    Returns: (patent_text, patent_id_stem)
    """
    if not HAS_REQUESTS:
        raise RuntimeError("requests/beautifulsoup4 필요: pip install requests beautifulsoup4")

    # Extract skey from URL for stem name
    skey_match = re.search(r'skey=(\d+)', url)
    stem = skey_match.group(1) if skey_match else "patent"

    headers = {
        "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36",
        "Accept-Language": "ko-KR,ko;q=0.9,en;q=0.8",
    }

    # 1) Main page (overview)
    print(f"  Fetching WIPS page: {url}")
    resp = requests.get(url, headers=headers, timeout=30)
    resp.encoding = 'utf-8'
    soup = BeautifulSoup(resp.text, 'html.parser')

    parts = []

    # Extract visible text sections
    for tag in soup.find_all(['div', 'td', 'span', 'p']):
        text = tag.get_text(strip=True)
        if len(text) > 30:  # Skip short navigation text
            parts.append(text)

    # 2) Try claims/description tab (different URL pattern)
    # WIPS uses AJAX tabs — try common endpoints
    base_url = url.split('?')[0]
    skey_param = url.split('?')[1] if '?' in url else ''

    # Try the claims iframe/tab
    for tab_type in ['DCI', 'AD']:  # DCI=claims, AD=full document
        try:
            tab_url = f"https://sd.wips.co.kr/wipslink/api/dusdshtm.wips?{skey_param}&tabType={tab_type}"
            r2 = requests.get(tab_url, headers=headers, timeout=20)
            r2.encoding = 'utf-8'
            s2 = BeautifulSoup(r2.text, 'html.parser')
            for tag in s2.find_all(['div', 'td', 'p', 'span']):
                text = tag.get_text(strip=True)
                if len(text) > 30:
                    parts.append(text)
        except Exception:
            continue

    # Deduplicate and join
    seen = set()
    unique_parts = []
    for p in parts:
        p_clean = p[:200]  # Use first 200 chars as key for dedup
        if p_clean not in seen:
            seen.add(p_clean)
            unique_parts.append(p)

    full_text = "\n\n".join(unique_parts)
    full_text = re.sub(r'\n{3,}', '\n\n', full_text).strip()

    if len(full_text) < 100:
        raise RuntimeError(f"WIPS 페이지에서 충분한 텍스트를 추출하지 못함 (len={len(full_text)})")

    print(f"  Extracted {len(full_text)} chars from WIPS page")
    return full_text, f"patent_{stem}"


# ============================================================
# LLM call (Chat Completions API)
# ============================================================
def call_llm(paper_text: str, model: str = DEFAULT_MODEL) -> Dict[str, Any]:
    api_key = os.getenv("OPENAI_API_KEY", "").strip()
    gemini_key = os.getenv("GEMINI_API_KEY", "").strip()

    # Truncate if too long
    text = paper_text[:110000]

    max_retries = 5
    backoff = 2.0

    # ── Gemini 모델 처리 ──────────────────────────────────
    if IS_GEMINI_MODEL(model):
        if not HAS_GEMINI:
            raise RuntimeError("google-generativeai 패키지가 설치되지 않았습니다. pip install google-generativeai")
        if not gemini_key:
            raise RuntimeError("GEMINI_API_KEY not set in .env")
        import google.generativeai as genai
        genai.configure(api_key=gemini_key)
        # response_mime_type \uc81c\uc678: gemini-2.5-preview \ud638\ud658\uc131 \ubb38\uc81c \ud68c\ud53c
        gmodel = genai.GenerativeModel(model)
        last_error = ""
        for attempt in range(max_retries):
            try:
                prompt = f"{SYSTEM_PROMPT}\n\n[\ubb38\uc11c \ud14d\uc2a4\ud2b8]\n{text}\n\n\ucd9c\ub825\uc740 \ubc18\ub4dc\uc2dc JSON\ub9cc \ucd9c\ub825. \ub9c8\ud06c\ub2e4\uc6b4 \ucf54\ub4dc\ube14\ub85d \uc5c6\uc774."
                resp = gmodel.generate_content(prompt)
                content = (resp.text or "").strip()
                if not content:
                    raise ValueError("Gemini returned empty response")
                json_match = re.search(r"```(?:json)?\s*([\s\S]+?)\s*```", content)
                if json_match:
                    content = json_match.group(1).strip()
                return json.loads(content)
            except Exception as e:
                last_error = f"{type(e).__name__}: {e}"
                print(f"  [WARN] Gemini attempt {attempt+1} failed: {last_error}")
                time.sleep(backoff * (2 ** attempt))

        # Gemini 실패 → OpenAI로 폴백
        if api_key:
            fallback_model = "gpt-5"
            print(f"  [FALLBACK] Gemini 실패. OpenAI({fallback_model})로 재시도합니다...")
            return call_llm(paper_text, model=fallback_model)
        raise RuntimeError(f"Gemini LLM call failed after retries. Last error: {last_error}")

    # ── OpenAI 모델 처리 ──────────────────────────────────
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY not set. .env 파일에 OPENAI_API_KEY를 추가하세요.")
    client = OpenAI(api_key=api_key)

    # Reasoning models (o1, o3, gpt-5 계열)
    is_reasoning_model = any(x in model.lower() for x in ["o1", "o3", "gpt-5", "gpt-o3"])

    last_error = ""
    truncated_text = text  # finish_reason=length 발생 시 자동 축소
    for attempt in range(max_retries):
        try:
            params: Dict[str, Any] = {
                "model": model,
                "messages": [
                    {"role": "system", "content": SYSTEM_PROMPT},
                    {"role": "user", "content": f"[문서 텍스트]\n{truncated_text}"},
                ],
            }
            if is_reasoning_model:
                # reasoning 모델은 내부 추론에 토큰을 많이 쓰므로 충분한 여유 필요
                params["max_completion_tokens"] = 32000
            else:
                params["response_format"] = {"type": "json_object"}
                params["max_completion_tokens"] = 4096
                params["temperature"] = 0.3
                params["top_p"] = 0.9

            resp = client.chat.completions.create(**params)
            content = (resp.choices[0].message.content or "").strip()
            finish_reason = resp.choices[0].finish_reason

            if not content:
                # finish_reason=length → 입력을 줄여서 출력 토큰 확보 후 재시도
                if finish_reason == "length" and len(truncated_text) > 20000:
                    new_len = max(20000, int(len(truncated_text) * 0.6))
                    print(f"  [INFO] finish_reason=length → 입력 {len(truncated_text)} → {new_len}자 축소 후 재시도")
                    truncated_text = truncated_text[:new_len]
                raise ValueError(f"LLM이 빈 응답을 반환했습니다. finish_reason={finish_reason}")

            json_match = re.search(r"```(?:json)?\s*([\s\S]+?)\s*```", content)
            if json_match:
                content = json_match.group(1).strip()

            try:
                return json.loads(content)
            except json.JSONDecodeError as je:
                # 응답 잘림(length) 의심 → 입력 축소 후 재시도
                if finish_reason == "length" and len(truncated_text) > 20000:
                    new_len = max(20000, int(len(truncated_text) * 0.6))
                    print(f"  [INFO] JSON 잘림 의심(finish_reason=length) → 입력 {len(truncated_text)} → {new_len}자 축소 후 재시도")
                    truncated_text = truncated_text[:new_len]
                raise ValueError(f"JSON 파싱 실패 ({je}). content head: {content[:120]!r}")
        except Exception as e:
            if "quota" in str(e).lower():
                raise
            last_error = f"{type(e).__name__}: {e}"
            print(f"  [WARN] LLM attempt {attempt+1} failed: {last_error}")
            time.sleep(backoff * (2 ** attempt))
    raise RuntimeError(f"LLM call failed after retries. Last error: {last_error}")


# ============================================================
# Post-processing (Stage C equivalent)
# ============================================================
@dataclass
class ProposedMethodItem:
    title: str
    details: List[str]

@dataclass
class Brief:
    pdf_stem: str
    doc_type: str
    title: str
    head_messages: List[str]
    purpose: List[str]
    prior_problems: List[str]
    proposed_method: List[ProposedMethodItem]
    improvements: List[str]
    conclusion: List[str]
    figure_captions_ko: List[str]
    paper_info: Dict[str, str]
    representative_figures: List[Dict]


def _smart_clean(text: str, max_chars: int, model: str = DEFAULT_MODEL) -> str:
    """글자수 초과 처리:
    - max_chars 이하: 그대로 반환
    - max_chars+1 ~ max_chars+10: 그대로 반환 (사람 검수)
    - max_chars+10 초과: API 재요약
    """
    t = _fix_units(_fix_punct(_strip_prefixes(text)))
    if len(t) <= max_chars + 10:
        return t
    print(f"  [INFO] 재요약 필요 ({len(t)}자 > {max_chars}+10): {t[:40]}...")
    return shorten_text_via_llm(t, max_chars, model=model)



def review_summary(raw: Dict[str, Any], model: str = DEFAULT_MODEL) -> Dict[str, Any]:
    """1차 요약 결과(JSON)를 검토자 관점에서 교정.

    검토 항목:
    1. 문장 자연스러움 (어색한 번역체, 불필요한 접속사)
    2. 오타 수정
    3. 지나치게 어려운 전문 용어 → 쉬운 표현 (업계 표준 약어는 유지)
    4. 숫자+단위 붙여쓰기 미준수 수정 (50 kHz → 50kHz 등)
    5. 지엽적 변수명 노출 검사 (논문 내부 기호 → 개념적 설명)
    6. 불릿 접두 기호, 문장 끝 마침표 등 형식 검사
    JSON 구조(키·배열 길이)는 절대 변경하지 않음.
    """
    api_key = os.getenv("OPENAI_API_KEY", "").strip()
    gemini_key = os.getenv("GEMINI_API_KEY", "").strip()

    # 검토할 필드만 추출 (이미지·메타 필드 제외)
    review_keys = ["title", "head_messages", "purpose", "prior_problems",
                   "proposed_method", "improvements"]
    subset = {k: raw.get(k) for k in review_keys if k in raw}

    REVIEW_PROMPT = f"""\
너는 자동차 기술 분야 요약문 검토 전문가다.
아래는 논문/특허 요약 결과 JSON이다. 원래의 공통 작성 규칙을 숙지한 상태에서 검토한다.

# 숙지할 원래 공통 작성 규칙 (검토 기준)
- 약어·전문 용어는 첫 등장 시 괄호 안에 영어 원문 또는 풀네임 병기
- 구체적인 알고리즘명·구조·프로세스 포함, 수식은 개념적으로 설명
- 지엽적 변수명·기호 (논문 내부에서만 쓰이는 Dp, phi12, VoH, P2ref 등) 사용 금지 → 개념적으로 풀어서 표현
- 업계 표준 약어(MOSFET, PWM, PID, SOC, BMS, IGBT, GNSS 등)는 유지
- 개조식(짧은 구/문장), 서술식 문단 금지
- 참조 번호(지지기35, 전축38 등) 절대 포함 금지
- 문장 끝 "." 금지
- 단위는 반드시 기호 사용: m, °, %
- 숫자와 단위 기호 사이 띄어쓰기 금지: 50kHz, 50kW, 1.2A, 3.5mm 등
- 불릿 접두 기호(•, -, ·, 번호) 절대 사용 금지
- 소수점은 18.4처럼 표기

# 검토 지침
1. 위 규칙을 위반한 부분만 교정
2. JSON 구조(키 이름, 배열 개수, 중첩 구조)는 절대 변경하지 말 것
3. proposed_method는 {{title: str, details: [str]}} 배열 형식을 유지
4. 내용을 추가·삭제하지 말고 표현만 다듬을 것
5. 이미 올바른 부분은 그대로 둘 것

입력 JSON:
{json.dumps(subset, ensure_ascii=False, indent=2)}

출력: 교정된 JSON만 출력 (같은 키 구조 유지, 마크다운 코드블록 없이)
"""

    try:
        if IS_GEMINI_MODEL(model) and HAS_GEMINI and gemini_key:
            import google.generativeai as genai
            genai.configure(api_key=gemini_key)
            gm = genai.GenerativeModel(model)
            resp = gm.generate_content(REVIEW_PROMPT + "\n\n\ucd9c\ub825\uc740 \ubc18\ub4dc\uc2dc JSON\ub9cc \ucd9c\ub825. \ub9c8\ud06c\ub2e4\uc6b4 \ucf54\ub4dc\ube14\ub85d \uc5c6\uc774.")
            raw = (resp.text or "").strip()
            jm = re.search(r"```(?:json)?\s*([\s\S]+?)\s*```", raw)
            corrected = json.loads(jm.group(1).strip() if jm else raw)
        elif api_key:
            client = OpenAI(api_key=api_key)
            is_reasoning = any(x in model.lower() for x in ["o1", "o3", "gpt-5", "gpt-o3"])
            params: Dict[str, Any] = {
                "model": model,
                "messages": [{"role": "user", "content": REVIEW_PROMPT}],
            }
            if is_reasoning:
                params["max_completion_tokens"] = 8000
            else:
                params["response_format"] = {"type": "json_object"}
                params["max_completion_tokens"] = 4096
                params["temperature"] = 0.2
            resp = client.chat.completions.create(**params)
            content = (resp.choices[0].message.content or "").strip()
            jm = re.search(r"```(?:json)?\s*([\s\S]+?)\s*```", content)
            corrected = json.loads(jm.group(1) if jm else content)
        else:
            print("  [WARN] review_summary: API key 없음. 검토 스킵.")
            return raw

        # 검토 키만 업데이트 (나머지 메타·이미지 키는 원본 유지)
        reviewed = dict(raw)
        for k in review_keys:
            if k in corrected:
                reviewed[k] = corrected[k]
        print("  [INFO] review_summary 완료.")
        return reviewed

    except Exception as e:
        print(f"  [WARN] review_summary 실패: {e}. 원본 사용.")
        return raw


def normalize_brief(raw: Dict[str, Any], pdf_stem: str, model: str = DEFAULT_MODEL) -> Brief:
    doc_type = raw.get("doc_type", "paper")
    if doc_type not in ("paper", "patent"):
        doc_type = "paper"

    is_patent = (doc_type == "patent")

    def _pc(text: str) -> str:
        """특허일 경우 object 번호 제거 후 반환."""
        if is_patent:
            text = _strip_patent_object_numbers(text)
        return text

    # Title: 50자 제한 (PPT 슬라이드 타이틀 — '■ ' 2자 포함시 52자)
    raw_title = _pc(raw.get("title", ""))
    title = _smart_clean(raw_title, 50, model)
    if not title:
        title = "제목 미추출"

    # Head messages: 150자 제한
    heads_raw = raw.get("head_messages") or []
    if not isinstance(heads_raw, list):
        heads_raw = []
    heads = []
    for h in heads_raw[:2]:
        h = _pc(str(h))
        heads.append(_smart_clean(h, 150, model))
    while len(heads) < 2:
        heads.append("요약 문장 미추출")

    # 기술 목적 / 문제점 / 개선효과 / 결론: 95자 제한
    def _clean_list(key: str, min_n: int, max_n: int, placeholder: str) -> List[str]:
        arr = raw.get(key) or []
        if not isinstance(arr, list):
            arr = []
        result = []
        for s in arr[:max_n]:
            s = _pc(str(s))
            cleaned = _smart_clean(s, 95, model)
            if cleaned:
                result.append(cleaned)
        while len(result) < min_n:
            result.append(placeholder)
        return result

    purpose = _clean_list("purpose", 2, 4, "기술 목적 미추출")
    prior_problems = _clean_list("prior_problems", 2, 4, "기존 문제점 미추출")
    improvements = _clean_list("improvements", 2, 4, "개선 효과 미추출")
    conclusion = _clean_list("conclusion", 2, 3, "결론 미추출")

    # proposed_method 3x3 — LLM이 array 또는 dict 형태로 반환할 수 있음
    pm_raw = raw.get("proposed_method") or []

    if isinstance(pm_raw, dict):
        converted = []
        for k, v in pm_raw.items():
            if isinstance(v, dict):
                converted.append({"title": k, "details": v.get("details") or []})
            elif isinstance(v, list):
                converted.append({"title": k, "details": v})
            else:
                converted.append({"title": k, "details": [str(v)]})
        pm_raw = converted

    if not isinstance(pm_raw, list):
        pm_raw = []

    pm_items = []
    for it in pm_raw[:4]:
        if not isinstance(it, dict):
            continue
        t = _pc(str(it.get("title", "")))
        t = _smart_clean(t, 75, model)
        ds = []
        for d in (it.get("details") or [])[:3]:
            d = _pc(str(d))
            d = _smart_clean(d, 90, model)
            if d:
                ds.append(d)
        while len(ds) < 3:
            ds.append("세부 내용 미추출")
        if not t:
            t = "기술 포인트 미추출"
        pm_items.append(ProposedMethodItem(title=t, details=ds))
    while len(pm_items) < 3:
        pm_items.append(ProposedMethodItem(title="기술 포인트 미추출", details=["세부 내용 미추출"] * 3))

    # 캡션: 60자 제한 (이미 vision 단계에서 생성, 빈 값 허용)
    caps_raw = raw.get("figure_captions_ko") or ["", ""]
    if not isinstance(caps_raw, list):
        caps_raw = ["", ""]

    def _is_english_caption(text: str) -> bool:
        """캡션이 영문인지 판별 — ASCII 비율이 70% 초과이면 영문으로 판단."""
        if not text:
            return False
        ascii_cnt = sum(1 for ch in text if ord(ch) < 128 and ch.strip())
        total = sum(1 for ch in text if ch.strip())
        return total > 0 and (ascii_cnt / total) > 0.70

    def _translate_single_caption(text: str, model: str) -> str:
        """단일 영문 캡션을 한국어로 번역. 실패 시 빈 문자열 반환."""
        try:
            results = _translate_captions_only(text, text, model=model)
            return results[0] if results else ""
        except Exception:
            return ""

    caps = []
    for c in caps_raw[:2]:
        c = str(c).strip()
        if c:
            # "Figure1.", "Fig.1.", "figure 1.", "Fig 3." 등 접두 패턴 제거
            c = re.sub(
                r'(?i)^\s*(?:figure|fig)\.?\s*\d+[a-z]?[\s.:\-\u2014]+',
                '', c
            ).strip()
            # 앞뒤 구두점 정리
            c = re.sub(r'^[\s.:\-]+', '', c).strip()
            # 영문 캡션 감지 → LLM 번역 시도
            if _is_english_caption(c):
                print(f"  [INFO] 영문 캡션 감지 → 한국어 번역 시도: '{c[:50]}'")
                c_translated = _translate_single_caption(c, model)
                if c_translated:
                    c = c_translated
                else:
                    # 번역 실패 시 빈 문자열 반환 (영문 원문 노출 차단)
                    print(f"  [WARN] 캡션 번역 실패 → 빈 캡션 처리")
                    caps.append("")
                    continue
            # "본 발명은/의/에" 등 금지 표현 제거
            c = _strip_invention_phrases(c)
            caps.append(_smart_clean(c, 60, model))
        else:
            caps.append("")  # 빈 값 허용
    while len(caps) < 2:
        caps.append("")

    pi = raw.get("paper_info") or {}
    if not isinstance(pi, dict):
        pi = {}
    for k in ["journal_or_patent_office", "paper_title", "institution", "doi_or_patent_no", "year", "month"]:
        pi[k] = str(pi.get(k) or "").strip()
    pi["institution"] = pick_primary_affiliation(pi.get("institution", ""))

    figs = raw.get("representative_figures") or []
    if not isinstance(figs, list):
        figs = []
    figs = figs[:2]

    return Brief(
        pdf_stem=pdf_stem, doc_type=doc_type, title=title,
        head_messages=heads, purpose=purpose, prior_problems=prior_problems,
        proposed_method=pm_items, improvements=improvements, conclusion=conclusion,
        figure_captions_ko=caps, paper_info=pi, representative_figures=figs,
    )


# ============================================================
# PPT helpers
# ============================================================
def _find_shape(slide, name: str):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def _remove_shape(shape):
    el = shape._element
    el.getparent().remove(el)


def _disable_auto_fit(tf):
    from pptx.oxml.ns import qn
    txBody = tf._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        for auto_el in bodyPr.findall(qn('a:spAutoFit')):
            bodyPr.remove(auto_el)
        for auto_el in bodyPr.findall(qn('a:normAutofit')):
            bodyPr.remove(auto_el)
        for auto_el in bodyPr.findall(qn('a:noAutofit')):
            bodyPr.remove(auto_el)
        if bodyPr.find(qn('a:normAutofit')) is None:
            from lxml import etree
            etree.SubElement(bodyPr, qn('a:normAutofit'))

def _set_paragraph_text(p, text: str):
    if not p.runs:
        p.add_run()
        p.runs[0].text = text
        return

    target_run = p.runs[0]
    max_len = -1
    for r in p.runs:
        if len(r.text.strip()) > max_len:
            max_len = len(r.text.strip())
            target_run = r

    for r in p.runs:
        if r != target_run:
            r.text = ""
            
    target_run.text = text

def _replace_text_keep_format(shape, text: str):
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs: return
    
    _set_paragraph_text(tf.paragraphs[0], text)
    
    # Remove extra paragraphs
    txBody = tf._txBody
    while len(tf.paragraphs) > 1:
        p_to_remove = tf.paragraphs[-1]._p
        txBody.remove(p_to_remove)

def _fill_multi_line_bullet(shape, items: list):
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs: return
    
    _disable_auto_fit(tf)
    
    txBody = tf._txBody
    from copy import deepcopy
    
    # duplicate last paragraph until length matches
    while len(tf.paragraphs) < len(items) and len(tf.paragraphs) > 0:
        new_p = deepcopy(tf.paragraphs[-1]._p)
        txBody.append(new_p)
        
    # remove extra paragraphs
    while len(tf.paragraphs) > len(items) and len(tf.paragraphs) > 1:
        txBody.remove(tf.paragraphs[-1]._p)
        
    for i, item in enumerate(items):
        if i < len(tf.paragraphs):
            _set_paragraph_text(tf.paragraphs[i], clean(item, 95))
            
    if not items and tf.paragraphs:
        _set_paragraph_text(tf.paragraphs[0], "")

def _fill_method_bullet(shape, item, idx: int):
    """(v1.2 호환용) method_bullet Shape에 제목 + 세부사항을 입력 (서식 완전 보존)"""
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs: return
    _disable_auto_fit(tf)
    
    txBody = tf._txBody
    from copy import deepcopy
    
    total_lines = 1 + min(len(item.details), 3)
    while len(tf.paragraphs) < total_lines and len(tf.paragraphs) > 0:
        new_p = deepcopy(tf.paragraphs[-1]._p)
        txBody.append(new_p)
        
    while len(tf.paragraphs) > total_lines and len(tf.paragraphs) > 1:
        txBody.remove(tf.paragraphs[-1]._p)
        
    circled_nums = ['①', '②', '③', '④']
    num = circled_nums[idx] if idx < len(circled_nums) else f"({idx+1})"
    
    if len(tf.paragraphs) > 0:
        _set_paragraph_text(tf.paragraphs[0], f"{num} {clean(item.title, 75)}")
    
    for i, d in enumerate(item.details[:3]):
        p_idx = i + 1
        if p_idx < len(tf.paragraphs):
            _set_paragraph_text(tf.paragraphs[p_idx], f"- {clean(d, 85)}")

# ============================================================
# Table cell helpers (v2 템플릿 — 표 기반)
# ============================================================
def _fill_table_cell(cell, items: list):
    """표 셀의 내용을 교체하되 템플릿 서식 보존."""
    tf = cell.text_frame
    if not tf.paragraphs: return
    
    txBody = tf._txBody
    from copy import deepcopy
    
    while len(tf.paragraphs) < len(items) and len(tf.paragraphs) > 0:
        new_p = deepcopy(tf.paragraphs[-1]._p)
        txBody.append(new_p)
        
    while len(tf.paragraphs) > len(items) and len(tf.paragraphs) > 1:
        txBody.remove(tf.paragraphs[-1]._p)
        
    for i, item in enumerate(items):
        if i < len(tf.paragraphs):
            _set_paragraph_text(tf.paragraphs[i], clean(item, 95))
            
    if not items and tf.paragraphs:
        _set_paragraph_text(tf.paragraphs[0], "")

def _fill_table_cell_method_title(cell, title_text: str):
    """method title 셀(5,7,9,11행)의 텍스트를 교체하되 서식 보존."""
    tf = cell.text_frame
    if not tf.paragraphs: return
    
    _set_paragraph_text(tf.paragraphs[0], clean(title_text, 75))
    
    txBody = tf._txBody
    while len(tf.paragraphs) > 1:
        p_to_remove = tf.paragraphs[-1]._p
        txBody.remove(p_to_remove)

def _fill_table_cell_method_details(cell, details: list):
    """method details 셀(6,8,10,12행)의 텍스트를 교체하되 서식 보존."""
    tf = cell.text_frame
    if not tf.paragraphs: return
    
    txBody = tf._txBody
    from copy import deepcopy
    
    while len(tf.paragraphs) < len(details) and len(tf.paragraphs) > 0:
        new_p = deepcopy(tf.paragraphs[-1]._p)
        txBody.append(new_p)
        
    while len(tf.paragraphs) > len(details) and len(tf.paragraphs) > 1:
        txBody.remove(tf.paragraphs[-1]._p)
        
    for i, d in enumerate(details[:3]):
        if i < len(tf.paragraphs):
            _set_paragraph_text(tf.paragraphs[i], f"- {clean(d, 85)}")
            
    if not details and tf.paragraphs:
        _set_paragraph_text(tf.paragraphs[0], "")

def _remove_table_row(tbl, row_idx: int):
    """표에서 특정 행을 XML 레벨로 제거."""
    tbl_element = tbl._tbl
    tr_elements = tbl_element.findall(qn('a:tr'))
    if 0 <= row_idx < len(tr_elements):
        tbl_element.remove(tr_elements[row_idx])


def _fill_method_table_cells(tbl, proposed_method: List['ProposedMethodItem']):
    """표의 method 영역(행 5~12)에 제안 기술 내용을 채움.
    method title: 행 5,7,9,11 / method details: 행 6,8,10,12
    채우기만 수행하고, 미사용 행 삭제는 별도로 호출.
    """
    n_methods = min(len(proposed_method), 4)
    circled_nums = ['①', '②', '③', '④']

    for i in range(n_methods):
        title_row = ROW_METHOD_TITLE_START + i * 2   # 5, 7, 9, 11
        detail_row = ROW_METHOD_DETAIL_START + i * 2  # 6, 8, 10, 12
        item = proposed_method[i]
        num = circled_nums[i] if i < len(circled_nums) else f"({i+1})"
        _fill_table_cell_method_title(tbl.cell(title_row, 0), f"{num} {item.title}")
        _fill_table_cell_method_details(tbl.cell(detail_row, 0), item.details)

    return n_methods


def _remove_unused_method_rows(tbl, n_methods: int):
    """method 항목이 4개 미만일 때 사용하지 않는 행을 표에서 삭제.
    반드시 모든 셀 채우기가 끝난 뒤에 호출할 것.
    """
    for i in range(3, n_methods - 1, -1):  # 역순: 3→2→...→n_methods
        detail_row = ROW_METHOD_DETAIL_START + i * 2  # 12, 10, ...
        title_row = ROW_METHOD_TITLE_START + i * 2    # 11, 9, ...
        _remove_table_row(tbl, detail_row)  # detail 먼저 삭제 (하위 행)
        _remove_table_row(tbl, title_row)   # title 삭제


# ============================================================
# Image replacement
# ============================================================
def _rects_intersect(a, b) -> bool:
    ax1, ay1, aw, ah = a
    bx1, by1, bw, bh = b
    return not (ax1 + aw <= bx1 or bx1 + bw <= ax1 or ay1 + ah <= by1 or by1 + bh <= ay1)



def replace_picture(slide, target_name: str, png_bytes: bytes, fallback_role: str = "concept"):
    sh = _find_shape(slide, target_name)
    if sh is None:
        # Fallback: find largest picture shapes on left side
        global SLIDE_W
        pics = []
        for s in slide.shapes:
            try:
                if s.shape_type not in (13, 14): continue
                if SLIDE_W and int(s.left) > int(SLIDE_W) * 0.55: continue
                area = int(s.width) * int(s.height)
                if area < 100000: continue
                pics.append((area, int(s.top), s))
            except: continue
        if len(pics) >= 2:
            pics.sort(key=lambda x: x[0], reverse=True)
            top2 = pics[:2]
            top2.sort(key=lambda x: x[1])
            sh = top2[0][2] if fallback_role == "concept" else top2[1][2]
        elif pics:
            sh = pics[0][2]
        else:
            return None

    left, top, width, height = int(sh.left), int(sh.top), int(sh.width), int(sh.height)
    try: _remove_shape(sh)
    except: pass

    # Remove overlapping large pictures
    for s in list(slide.shapes):
        try:
            if s.shape_type not in (13, 14): continue
            a = int(s.width) * int(s.height)
            if a < int(width) * int(height) * 0.25: continue
            r = (int(s.left), int(s.top), int(s.width), int(s.height))
            if _rects_intersect(r, (left, top, width, height)):
                _remove_shape(s)
        except: continue

    bio = io.BytesIO(png_bytes)
    slide.shapes.add_picture(bio, left, top, width=width, height=height)
    return left, top, width, height


# ============================================================
# Slide duplication
# ============================================================
def duplicate_slide(prs: Presentation, slide_index: int = 0):
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    source = prs.slides[slide_index]
    new_slide = prs.slides.add_slide(source.slide_layout)
    for sh in list(new_slide.shapes):
        _remove_shape(sh)

    for sh in source.shapes:
        new_el = deepcopy(sh._element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    rid_map = {}
    for rId, rel in source.part.rels.items():
        try:
            if rel.reltype != RT.IMAGE: continue
            blob = rel.target_part.blob
            new_part = prs.part.package.get_or_add_image_part(io.BytesIO(blob))
            new_rId = new_slide.part.relate_to(new_part, RT.IMAGE)
            rid_map[rId] = new_rId
        except: continue

    if rid_map:
        blips = new_slide._element.xpath(".//a:blip")
        for blip in blips:
            try:
                old = blip.get(qn("r:embed"))
                if old and old in rid_map:
                    blip.set(qn("r:embed"), rid_map[old])
            except: continue
    return new_slide


# ============================================================
# Build source text
# ============================================================
def build_source_text(brief: Brief) -> str:
    pi = brief.paper_info
    doi_or_no = pi.get("doi_or_patent_no", "").strip() or "N/A"
    inst = pi.get("institution", "").strip() or "N/A"
    year = pi.get("year", "").strip() or "N/A"
    if brief.doc_type == "patent":
        return f"*출처: {doi_or_no} ({inst}, {year})"
    else:
        return f"*출처: {doi_or_no} ({inst}, {year})"


# ============================================================
# Fill one slide (v2 — 개별 Shape 직접 매핑)
# ============================================================
def fill_slide(slide, brief: Brief, doc_idx: int, concept_fig: Optional[Dict] = None,
               effect_fig: Optional[Dict] = None):
    # 1) number_info: 문서 라벨
    sh_label = _find_shape(slide, SHAPE_NUMBER_INFO)
    if sh_label:
        kr_type = "특허" if brief.doc_type == "patent" else "논문"
        _replace_text_keep_format(sh_label, f"{kr_type}_후보기술 개요서_{doc_idx}")

    # 2) title: 제목
    sh_title = _find_shape(slide, SHAPE_TITLE)
    if sh_title:
        title_text = f"{brief.title}".strip()
        if not title_text.startswith("■"):
            title_text = f"■ {title_text}"
        _replace_text_keep_format(sh_title, title_text)

    # 3) summary_point_1-2: 헤드메시지 (2문장)
    sh_summary = _find_shape(slide, SHAPE_SUMMARY)
    if sh_summary:
        _fill_multi_line_bullet(sh_summary, brief.head_messages[:2])

    # 4~7) 표 기반: 기술목적 / 문제점 / 제안기술 / 개선효과
    sh_table = _find_shape(slide, TABLE_NAME)
    if sh_table and sh_table.has_table:
        tbl = sh_table.table
        # 기술 목적 (행 1)
        _fill_table_cell(tbl.cell(ROW_PURPOSE_CONTENT, 0), brief.purpose[:4])
        # 기존 기술의 문제점 (행 3)
        _fill_table_cell(tbl.cell(ROW_PROBLEMS_CONTENT, 0), brief.prior_problems[:4])
        # 제안 기술 (행 5~12) — 채우기만
        n_methods = _fill_method_table_cells(tbl, brief.proposed_method)
        # 개선 효과 (행 14) — 행 삭제 전에 먼저 채우기
        _fill_table_cell(tbl.cell(ROW_EFFECT_CONTENT, 0), brief.improvements[:4])
        # 미사용 method 행 삭제 (모든 셀 채우기 완료 후)
        _remove_unused_method_rows(tbl, n_methods)

    # 8) source_info: 출처
    sh_src = _find_shape(slide, SHAPE_SOURCE)
    if sh_src:
        _replace_text_keep_format(sh_src, build_source_text(brief))

    # 9) Figures — diagram1 (컨셉), diagram2 (결과)
    cap1 = brief.figure_captions_ko[0] if brief.figure_captions_ko else ""
    cap2 = brief.figure_captions_ko[1] if len(brief.figure_captions_ko) > 1 else ""

    # representative_figures에서 fig_number 추출 (이미지 없을 때 캡션에 사용)
    fig1_number = ""
    fig2_number = ""
    if brief.representative_figures:
        if len(brief.representative_figures) > 0:
            fig1_number = brief.representative_figures[0].get("fig_number", "")
        if len(brief.representative_figures) > 1:
            fig2_number = brief.representative_figures[1].get("fig_number", "")

    def _build_caption_label(fig_dict: Optional[Dict], cap_ko: str,
                             fig_number_from_brief: str, slot_idx: int) -> str:
        """이미지 슬롯 1개의 표시용 캡션 텍스트를 생성.

        우선순위 (요청서 §11):
          1) translated_caption (한국어 번역본) — figure dict 또는 cap_ko 인자
          2) original_caption (PDF 원문 영문 캡션) — figure dict
          3) fallback: "논문 p.{page} 대표 이미지 - 캡션 확인 필요"

        부가 규칙 (요청서 §10):
          - 캡션이 이미 'Fig./Figure/그림/도면'으로 시작하면 prefix를 다시 붙이지 않음
          - normalize_caption()으로 'Fig.Fig.' 같은 중복 prefix 제거
          - "그림1", "그림2" 단독 라벨은 출력하지 않음
        """
        # fig_number 결정 우선순위: 이미지 객체의 fig_label/figure_number → brief.representative_figures
        fig_num = ""
        if fig_dict:
            fig_num = (_normalize_fig_number(fig_dict.get("fig_label", ""))
                       or _normalize_fig_number(fig_dict.get("figure_number", ""))
                       or "")
        if not fig_num:
            fig_num = _normalize_fig_number(fig_number_from_brief) or ""

        # 1) translated_caption — cap_ko 인자(=brief.figure_captions_ko) 우선, 그 다음 figure dict
        translated = (cap_ko or "").strip()
        if not translated and fig_dict:
            translated = (fig_dict.get("translated_caption") or "").strip()

        # 2) original_caption — figure dict
        original = ""
        if fig_dict:
            original = (fig_dict.get("original_caption")
                        or fig_dict.get("nearest_caption")
                        or "").strip()

        chosen = translated or original
        chosen_source = "translated_caption" if translated else ("original_caption" if original else "")

        if chosen:
            chosen = normalize_caption(chosen)
            # 단독 라벨("그림1", "Figure 1" 등)은 캡션으로 인정하지 않음
            if re.match(r"^(?:Fig\.?|Figure|그림|도면)\s*\d+\s*$", chosen, re.IGNORECASE):
                chosen = ""
            elif re.match(r"^\s*(?:Fig\.?|Figure|그림|도면)", chosen, re.IGNORECASE):
                # 이미 Fig./Figure/그림 등으로 시작 → prefix 재첨부 금지
                return chosen
            else:
                if fig_num:
                    return f"Fig. {fig_num}. {chosen}"
                return chosen

        # caption 미검출 fallback (요청서 §12)
        page_str = ""
        if fig_dict:
            page_str = str(fig_dict.get("page") or fig_dict.get("page_number") or "")
        print(f"  [WARN] caption fallback used: slot={slot_idx} page={page_str or '?'} "
              f"fig_number={fig_num or '?'} (no translated/original caption available)")
        if page_str:
            return f"논문 p.{page_str} 대표 이미지 - 캡션 확인 필요"
        return "대표 이미지 - 캡션 확인 필요"

    def _build_image_missing_label(fig_number_from_brief: str, slot_idx: int) -> str:
        """이미지 추출 fallback 표시(요청서 §12). 'Fig.Fig.' 중복 prefix 방지."""
        fn = _normalize_fig_number(fig_number_from_brief) or ""
        ref = f"Fig. {fn}. " if fn else ""
        msg = normalize_caption(f"{ref}이미지 추출 불가로 원본 확인 필요")
        print(f"  [ERROR] image fallback used: slot={slot_idx} fig_number={fn or '?'}")
        return msg

    if concept_fig and concept_fig.get("png_bytes"):
        box = replace_picture(slide, SHAPE_DIAGRAM1, concept_fig["png_bytes"], "concept")
        if box:
            sh_cap1 = _find_shape(slide, SHAPE_CAPTION1)
            if sh_cap1:
                cap_text = _build_caption_label(concept_fig, cap1, fig1_number, 1)
                _replace_text_keep_format(sh_cap1, clean(cap_text, 80, hard_cut=True))
    else:
        # 이미지 추출 불가 → 캡션 영역에 안내 텍스트만 기입, 이미지 영역은 그대로 유지
        sh_cap1 = _find_shape(slide, SHAPE_CAPTION1)
        if sh_cap1:
            _replace_text_keep_format(sh_cap1, _build_image_missing_label(fig1_number, 1))

    if effect_fig and effect_fig.get("png_bytes"):
        box = replace_picture(slide, SHAPE_DIAGRAM2, effect_fig["png_bytes"], "result")
        if box:
            sh_cap2 = _find_shape(slide, SHAPE_CAPTION2)
            if sh_cap2:
                cap_text = _build_caption_label(effect_fig, cap2, fig2_number, 2)
                _replace_text_keep_format(sh_cap2, clean(cap_text, 80, hard_cut=True))
    else:
        # 이미지 추출 불가 → 캡션 영역에 안내 텍스트만 기입, 이미지 영역은 그대로 유지
        sh_cap2 = _find_shape(slide, SHAPE_CAPTION2)
        if sh_cap2:
            _replace_text_keep_format(sh_cap2, _build_image_missing_label(fig2_number, 2))



# ============================================================
# MD summary output
# ============================================================
def build_md(brief: Brief) -> str:
    lines = []
    lines.append(f"# {brief.pdf_stem}. {brief.title}")
    lines.append("")
    lines.append("## 헤드메시지")
    for h in brief.head_messages:
        lines.append(f"- {h}")
    lines.append("")
    lines.append("## 기술목적")
    for x in brief.purpose: lines.append(f"- {x}")
    lines.append("")
    lines.append("## 기존 기술의 문제점")
    for x in brief.prior_problems: lines.append(f"- {x}")
    lines.append("")
    lines.append("## 제안 기술의 구체적인 방법 혹은 기술컨셉")
    for i, it in enumerate(brief.proposed_method, 1):
        lines.append(f"{i}) {it.title}")
        for d in it.details:
            lines.append(f"   - {d}")
        lines.append("")
    lines.append("## 개선효과")
    for x in brief.improvements: lines.append(f"- {x}")
    lines.append("")

    pi = brief.paper_info
    lines.append("## 논문/특허 정보")
    lines.append(f"- Journal/Patent Office: {pi.get('journal_or_patent_office', '')}")
    lines.append(f"- Title: {pi.get('paper_title', '')}")
    lines.append(f"- Institution: {pi.get('institution', '')}")
    lines.append(f"- DOI/Patent No: {pi.get('doi_or_patent_no', '')}")
    lines.append(f"- Year: {pi.get('year', '')} / Month: {pi.get('month', '')}")
    lines.append("")

    if brief.representative_figures:
        lines.append("## 대표 이미지")
        for fig in brief.representative_figures:
            lines.append(f"- Page {fig.get('page','?')}, {fig.get('fig_number','?')}: "
                         f"{fig.get('fig_title','?')} (role: {fig.get('role','?')})")
    lines.append("")
    lines.append("## Figure 캡션")
    for c in brief.figure_captions_ko:
        lines.append(f"- {c}")
    lines.append("")
    lines.append(build_source_text(brief))
    return "\n".join(lines)


# ============================================================
# Main
# ============================================================
def process_one_document(
    text: str,
    stem: str,
    model: str,
    save_json: bool,
    save_md: bool,
    out_dir: Path,
    concept_fig: Optional[Dict] = None,
    effect_fig: Optional[Dict] = None,
    vision_caps: Optional[List[str]] = None,
    precomputed_raw: Optional[Dict] = None,
    source_pdf: str = "",
) -> Tuple[Brief, Optional[Dict], Optional[Dict]]:
    """LLM 호출 → 정규화 → JSON/MD 저장, Brief 반환.
    precomputed_raw가 주어지면 LLM 재호출을 스킵하고 해당 결과를 사용."""
    if precomputed_raw is not None:
        raw = precomputed_raw
        print(f"  Using pre-computed LLM result (skipping call_llm).")
    else:
        print(f"  Calling LLM ({model})...")
        raw = call_llm(text, model=model)

    if vision_caps and len(vision_caps) == 2:
        raw["figure_captions_ko"] = vision_caps

    brief = normalize_brief(raw, stem, model=model)

    if save_json:
        jpath = out_dir / f"{stem}.json"
        # 원문 파일명 메타정보 삽입 (비어있으면 스킵)
        if source_pdf:
            raw["source_pdf"] = source_pdf
        with open(jpath, "w", encoding="utf-8") as f:
            json.dump(raw, f, ensure_ascii=False, indent=2)
        print(f"  Saved JSON: {jpath}")

    if save_md:
        mpath = out_dir / f"{stem}.md"
        with open(mpath, "w", encoding="utf-8") as f:
            f.write(build_md(brief))
        print(f"  Saved MD: {mpath}")


    return brief, concept_fig, effect_fig


def main():
    global SLIDE_W, SLIDE_H

    ap = argparse.ArgumentParser(description="후보기술개요서 자동 요약 봇")
    ap.add_argument("--pdf_dir", default="", help="PDF 파일이 있는 디렉토리")
    ap.add_argument("--out_dir", required=True, help="출력 디렉토리")
    ap.add_argument("--template", required=True, help="PPT 템플릿 파일 경로")
    ap.add_argument("--wips_urls", default="", help="WIPS 특허 URL (쉼표 구분, 여러 개 가능)")
    ap.add_argument("--model", default=DEFAULT_MODEL, help=f"LLM 모델 (default: {DEFAULT_MODEL})")
    ap.add_argument("--max_pages", type=int, default=20, help="PDF에서 읽을 최대 페이지 수")
    ap.add_argument("--out_name", default="output_briefs.pptx", help="출력 PPT 파일명")
    ap.add_argument("--save_json", action="store_true", help="JSON 중간결과 저장")
    ap.add_argument("--save_md", action="store_true", help="MD 요약 저장")
    ap.add_argument("--no_images", action="store_true", help="이미지 추출/삽입 스킵")
    args = ap.parse_args()

    out_dir = Path(args.out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    if not os.path.exists(args.template):
        raise RuntimeError(f"Template not found: {args.template}")

    # Collect jobs: list of (source_type, source_data)
    jobs: List[Tuple[str, Any]] = []  # ("pdf", Path) or ("wips", url_str)

    # PDFs
    if args.pdf_dir and os.path.isdir(args.pdf_dir):
        pdfs = sorted(Path(args.pdf_dir).glob("*.pdf"), key=lambda p: p.name.lower())
        for p in pdfs:
            jobs.append(("pdf", p))

    # WIPS URLs
    if args.wips_urls:
        urls = [u.strip() for u in args.wips_urls.split(",") if u.strip()]
        for u in urls:
            jobs.append(("wips", u))

    if not jobs:
        raise RuntimeError("처리할 문서가 없습니다. --pdf_dir 또는 --wips_urls를 지정하세요.")

    prs = Presentation(args.template)
    SLIDE_W = int(prs.slide_width)
    SLIDE_H = int(prs.slide_height)

    # Slide 0 = 첫 번째 콘텐츠 템플릿 (새 템플릿 기준)
    NUM_INTRO_SLIDES = 0
    base_slide_idx = NUM_INTRO_SLIDES

    if len(prs.slides) <= base_slide_idx:
        raise RuntimeError(
            f"템플릿 슬라이드 수({len(prs.slides)})가 부족합니다. "
        )

    for idx, (src_type, src_data) in enumerate(jobs):
        concept_fig, effect_fig, vision_caps = None, None, None
        precomputed_raw = None

        if src_type == "pdf":
            pdf_path = src_data
            stem = pdf_path.stem
            print(f"\n== [{idx+1}/{len(jobs)}] Processing PDF: {pdf_path.name}")
            _t_total = time.time()

            _t0 = time.time()
            with fitz.open(str(pdf_path)) as doc:
                _t0 = _timer_log("pdf_open", _t0)

                text = extract_text(doc, max_pages=args.max_pages)
                _t0 = _timer_log("text_parse", _t0)

                # Step 1: 텍스트 LLM을 먼저 호출하여 representative_figures 확보
                print(f"  Calling LLM ({args.model})...")
                precomputed_raw = call_llm(text, model=args.model)
                _t0 = _timer_log("ai_summarization", _t0)

                # 논문만 fig_title 전달 — 특허는 캡션이 없어 텍스트 LLM 선정이 부정확하므로 Vision LLM이 독립 판단
                _is_patent = precomputed_raw.get("doc_type") == "patent"
                rep_figs = [] if _is_patent else (precomputed_raw.get("representative_figures") or [])

                # Step 2: Vision LLM에 representative_figures의 fig_title을 전달
                if not args.no_images:
                    try:
                        candidates = extract_images(doc)
                        _t0 = _timer_log("figure_candidate_build", _t0)

                        concept_fig, effect_fig, vision_caps = pick_two_figures_with_vision(
                            candidates, text, model=args.model,
                            representative_figures=rep_figs,
                        )
                        _t0 = _timer_log("ai_selection_and_translation", _t0)
                    except Exception as e:
                        print(f"  [WARN] Image extraction failed: {e}")
                        vision_caps = ["", ""]

        elif src_type == "wips":
            url = src_data
            print(f"\n== [{idx+1}/{len(jobs)}] Processing WIPS URL: {url}")
            _t_total = time.time()
            text, stem = scrape_wips_patent(url)

        else:
            continue

        # 원문 파일명 결정 (PDF는 파일명, WIPS는 URL)
        _source_pdf = pdf_path.name if src_type == "pdf" else (src_data if src_type == "wips" else "")

        _t_ppt = time.time()
        brief, concept_fig, effect_fig = process_one_document(
            text=text, stem=stem, model=args.model,
            save_json=args.save_json, save_md=args.save_md,
            out_dir=out_dir, concept_fig=concept_fig, effect_fig=effect_fig,
            vision_caps=vision_caps, precomputed_raw=precomputed_raw,
            source_pdf=_source_pdf,
        )

        # 매번 원본 템플릿(base_slide_idx)을 복제하여 새 슬라이드 생성
        # → method 행 삭제 등으로 원본이 변형되지 않도록 보장
        slide = duplicate_slide(prs, slide_index=base_slide_idx)

        fill_slide(slide, brief, doc_idx=idx+1, concept_fig=concept_fig, effect_fig=effect_fig)
        _timer_log("ppt_generation", _t_ppt)
        _timer_log("total", _t_total)
        print(f"  → Slide {idx+1} filled: {brief.title}")

    # 원본 템플릿 슬라이드(base_slide_idx) 제거 — 복제본만 남김
    from lxml import etree as _et_main
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    slide_list = prs.slides._sldIdLst
    # base_slide_idx 이하의 슬라이드(원본 템플릿 + 지침 슬라이드) 제거
    n_remove = base_slide_idx + 1
    slide_ids_to_remove = list(slide_list)[:n_remove]
    for sldId in slide_ids_to_remove:
        rId = sldId.get("r:id")
        if rId:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass
        slide_list.remove(sldId)

    # 고아 슬라이드 관계 정리 (ZIP 중복 방지)
    active_rids = set()
    for sldId in slide_list:
        rid = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id") or sldId.get("r:id")
        if rid:
            active_rids.add(rid)
    for rId, rel in list(prs.part.rels.items()):
        if rel.reltype == RT.SLIDE and rId not in active_rids:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass

    out_path = str(out_dir / args.out_name)
    prs.save(out_path)
    print(f"\n[DONE] Saved: {out_path}  ({len(jobs)} slide(s))")

    if Image is None:
        print("[WARN] Pillow 미설치 — 이미지 여백 트리밍 비활성화됨 (pip install pillow)")


if __name__ == "__main__":
    main()
