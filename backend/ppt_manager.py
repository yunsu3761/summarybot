# ppt_manager.py
# -------------------------------------------------------
# PPT 슬라이드 순서 재배치 및 레이블 업데이트 모듈
# -------------------------------------------------------

import io
import re
from copy import deepcopy
from pathlib import Path
from typing import List, Dict, Optional, Tuple

from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn

# 슬라이드 내 문서 레이블 shape 이름 (build_briefs_v2.py의 SHAPE_NUMBER_INFO와 동일)
SHAPE_DOC_LABEL = "number_info"


# ============================================================
# Shape 유틸리티
# ============================================================
def _find_shape(slide, name: str):
    """슬라이드에서 이름으로 shape을 찾습니다."""
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def _remove_shape(shape):
    """슬라이드에서 shape을 제거합니다."""
    el = shape._element
    el.getparent().remove(el)


def _extract_style(run) -> dict:
    """run에서 폰트 스타일을 추출합니다."""
    return {
        'name': run.font.name,
        'size': run.font.size,
        'bold': run.font.bold,
    }


def _set_run_style(run, style: dict):
    """run에 폰트 스타일을 적용합니다."""
    if not style:
        return
    if style.get('name'):
        run.font.name = style['name']
        try:
            rPr = run._r.get_or_add_rPr()
            rPr.set(qn("a:ea"), style['name'])
        except Exception:
            pass
    if style.get('size'):
        run.font.size = style['size']
    if style.get('bold') is not None:
        run.font.bold = style['bold']


def _set_paragraph_text(p, text: str):
    """단일 paragraph의 텍스트를 교체하면서 기존 run 포맷을 유지합니다."""
    if not p.runs:
        p.add_run()
        p.runs[0].text = text
        return

    target_run = p.runs[0]
    max_len = -1
    for r in p.runs:
        if len(r.text.strip()) > max_len:
            max_len = len(r.text.strip())
            target_run = r

    for r in p.runs:
        if r != target_run:
            r.text = ""

    target_run.text = text


def _replace_text_keep_format(shape, text: str):
    """shape의 텍스트를 교체하면서 기존 포맷을 유지합니다.

    tf.clear()를 사용하지 않고, 기존 paragraph/run 구조를 보존하면서
    텍스트만 교체합니다. 이는 python-pptx 저장 시 ZIP 충돌을 방지합니다.
    """
    if not shape or not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        return

    _set_paragraph_text(tf.paragraphs[0], text)

    # Remove extra paragraphs
    txBody = tf._txBody
    while len(tf.paragraphs) > 1:
        p_to_remove = tf.paragraphs[-1]._p
        txBody.remove(p_to_remove)


# ============================================================
# 슬라이드 레이블에서 번호 추출
# ============================================================
_LABEL_NUMBER_PATTERN = re.compile(r'(\d+)')


def extract_slide_number(slide) -> Optional[int]:
    """슬라이드의 문서 레이블에서 번호를 추출합니다."""
    sh = _find_shape(slide, SHAPE_DOC_LABEL)
    if not sh or not sh.has_text_frame:
        return None
    text = sh.text_frame.text.strip()
    match = _LABEL_NUMBER_PATTERN.search(text)
    if match:
        return int(match.group(1))
    return None


# ============================================================
# 슬라이드 순서 재배치
# ============================================================
def reorder_slides(input_pptx: str, new_order: List[int], output_pptx: str):
    """
    PPT 슬라이드의 물리적 순서를 재배치합니다.
    """
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    prs = Presentation(input_pptx)
    
    # 0. 고아 슬라이드 파트 정리 (ZIP 중복 방지)
    sld_id_lst = prs.slides._sldIdLst
    active_rids = set()
    for sldId in sld_id_lst:
        rid = sldId.get(qn("r:id"))
        if rid:
            active_rids.add(rid)

    for rId, rel in list(prs.part.rels.items()):
        if rel.reltype == RT.SLIDE and rId not in active_rids:
            try:
                prs.part.drop_rel(rId)
            except Exception:
                pass

    # 1. 기존 슬라이드 ID 리스트를 복사하여 보관 (SnapShot)
    current_slides = list(sld_id_lst) 

    # 2. 유효성 검증
    num_slides = len(current_slides)
    if len(new_order) != num_slides:
        raise ValueError(f"입력된 순서 리스트({len(new_order)})와 실제 슬라이드 수({num_slides})가 일치해야 합니다.")

    for idx in new_order:
        if idx < 0 or idx >= num_slides:
            raise ValueError(f"잘못된 슬라이드 인덱스: {idx}")

    # 3. 부모 요소에서 모든 자식(슬라이드 ID)을 제거합니다.
    for sld in current_slides:
        sld_id_lst.remove(sld)

    # 4. 새로운 순서에 맞춰 하나씩 다시 추가 (붙여넣기)
    for new_idx in new_order:
        sld_id_lst.append(current_slides[new_idx])

    # 5. 저장
    prs.save(output_pptx)
    return len(new_order)


def reorder_by_number_mapping(input_pptx: str, mapping: Dict[int, int],
                               output_pptx: str) -> int:
    """
    번호 변환 매핑에 따라 슬라이드를 재배치합니다.

    입력 PPT의 각 슬라이드는 before 번호 순서로 배치되어 있으며,
    레이블에 "논문(또는 특허)_후보기술_개요서_{before번호}"가 표기되어 있습니다.

    처리 순서:
      1. 고아 슬라이드 파트 정리 (ZIP 중복 방지)
      2. 각 슬라이드의 레이블에서 before 번호를 추출
      3. mapping({before번호: after번호})을 적용하여 after 번호 결정
      4. 레이블 텍스트의 번호를 after 번호로 변경
      5. after 번호 오름차순으로 슬라이드를 물리적 재배치
      6. 한 번에 저장 (형식 보존)

    Args:
        input_pptx: 입력 PPTX 파일 경로
        mapping: {before번호: after번호}
        output_pptx: 출력 PPTX 파일 경로

    Returns:
        재배치된 슬라이드 수
    """
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    prs = Presentation(input_pptx)

    # ── 0단계: 고아 슬라이드 파트 정리 ──
    # sldIdLst에서 참조하는 rId 목록 수집
    sld_id_lst = prs.slides._sldIdLst
    active_rids = set()
    for sldId in sld_id_lst:
        rid = sldId.get(qn("r:id"))
        if rid:
            active_rids.add(rid)

    # 프레젠테이션 파트의 모든 slide relationship 중
    # sldIdLst에서 참조하지 않는 것을 제거 (고아 파트)
    orphan_rids = []
    for rId, rel in list(prs.part.rels.items()):
        if rel.reltype == RT.SLIDE and rId not in active_rids:
            orphan_rids.append(rId)

    for rId in orphan_rids:
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass

    slides = list(prs.slides)

    # ── 1단계: 각 슬라이드의 before 번호 추출 ──
    slide_info = []  # [(slide_index, before_num)]
    for i, slide in enumerate(slides):
        num = extract_slide_number(slide)
        if num is None:
            num = i + 1  # 번호를 찾지 못하면 순서번호 사용
        slide_info.append((i, num))

    # ── 2단계: mapping 적용 → (slide_index, before_num, after_num) ──
    slide_mapping = []
    for idx, before_num in slide_info:
        after_num = mapping.get(before_num, before_num)  # 매핑 없으면 기존 번호 유지
        slide_mapping.append((idx, before_num, after_num))

    # ── 3단계: 레이블 번호를 after 번호로 변경 (물리적 재배치 전, 원본 인덱스로 접근) ──
    # 숫자 또는 '#' 패턴 모두 교체 (v3 템플릿의 '#' 대응)
    _label_replace_pattern = re.compile(r'(\d+|#)')
    for idx, before_num, after_num in slide_mapping:
        slide = slides[idx]
        sh = _find_shape(slide, SHAPE_DOC_LABEL)
        if sh and sh.has_text_frame:
            text = sh.text_frame.text.strip()
            new_text = _label_replace_pattern.sub(str(after_num), text, count=1)
            _replace_text_keep_format(sh, new_text)

    # ── 4단계: after 번호 오름차순으로 물리적 순서 재배치 (XML sldIdLst 조작) ──
    slide_mapping.sort(key=lambda x: x[2])  # after 번호 기준 정렬
    new_order = [item[0] for item in slide_mapping]

    current_slides = list(sld_id_lst)

    # 기존 슬라이드 ID를 모두 제거
    for sld in current_slides:
        sld_id_lst.remove(sld)

    # 새 순서대로 다시 추가
    for new_idx in new_order:
        sld_id_lst.append(current_slides[new_idx])

    # ── 5단계: 한 번에 저장 ──
    prs.save(output_pptx)
    return len(new_order)


# ============================================================
# 슬라이드 레이블 업데이트
# ============================================================
def update_slide_labels(pptx_path: str, label_map: Dict[int, str],
                        output_path: str = None):
    """
    슬라이드 내 문서 레이블 텍스트를 업데이트합니다.

    Args:
        pptx_path: PPTX 파일 경로
        label_map: {슬라이드 인덱스(0-based): 새 레이블 텍스트}
        output_path: 출력 경로 (없으면 원본 덮어쓰기)
    """
    prs = Presentation(pptx_path)

    for i, slide in enumerate(prs.slides):
        if i in label_map:
            sh = _find_shape(slide, SHAPE_DOC_LABEL)
            if sh:
                _replace_text_keep_format(sh, label_map[i])

    save_path = output_path or pptx_path
    prs.save(save_path)


def update_slide_numbers(pptx_path: str, new_numbers: List[int],
                         doc_types: List[str] = None,
                         output_path: str = None):
    """
    슬라이드 레이블의 번호만 업데이트합니다.

    레이블 형식: "{논문/특허}_후보기술 개요서_{번호}"

    Args:
        pptx_path: PPTX 파일 경로
        new_numbers: 각 슬라이드의 새 번호 리스트
        doc_types: 각 슬라이드의 문서 유형 (paper/patent)
        output_path: 출력 경로
    """
    prs = Presentation(pptx_path)

    for i, slide in enumerate(prs.slides):
        if i < len(new_numbers):
            sh = _find_shape(slide, SHAPE_DOC_LABEL)
            if sh:
                doc_type = "논문"
                if doc_types and i < len(doc_types):
                    doc_type = "특허" if doc_types[i] == "patent" else "논문"
                label = f"{doc_type}_후보기술 개요서_{new_numbers[i]}"
                _replace_text_keep_format(sh, label)

    save_path = output_path or pptx_path
    prs.save(save_path)


# ============================================================
# PPT 정보 조회
# ============================================================
def get_pptx_info(pptx_path: str) -> List[Dict]:
    """
    PPTX 파일의 슬라이드 정보를 반환합니다.

    Returns:
        List[Dict]: 각 슬라이드의 {index, number, label, shapes}
    """
    prs = Presentation(pptx_path)
    info = []

    for i, slide in enumerate(prs.slides):
        slide_info = {
            "index": i,
            "number": extract_slide_number(slide),
            "label": "",
            "shapes": [],
        }

        # 레이블 텍스트
        sh = _find_shape(slide, SHAPE_DOC_LABEL)
        if sh and sh.has_text_frame:
            slide_info["label"] = sh.text_frame.text.strip()

        # 주요 shape 이름들
        for sh in slide.shapes:
            slide_info["shapes"].append(sh.name)

        info.append(slide_info)

    return info
